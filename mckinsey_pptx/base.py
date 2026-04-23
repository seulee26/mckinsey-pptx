"""Common slide primitives: title, underline, footer, section marker, text/shape helpers."""
from __future__ import annotations
from typing import Optional, Iterable

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from .theme import Theme, DEFAULT_THEME


# ---------- text helpers ----------

def set_run(run, text, *, size=None, bold=False, italic=False, color=None,
            family=None):
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if family:
        run.font.name = family
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left_in, top_in, width_in, height_in, *, fill=None,
                line=None, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line is None:
        tb.line.fill.background()
    return tb


def write_paragraph(tf, text, *, size, bold=False, italic=False, color=None,
                    family="Arial", align=PP_ALIGN.LEFT, space_before=None,
                    space_after=None, level=0, bullet=False, first=False):
    if first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, italic=italic, color=color,
            family=family)
    if bullet:
        _set_bullet(p, color)
    else:
        _clear_bullet(p)
    return p


def _set_bullet(paragraph, color: Optional[RGBColor]):
    pPr = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu = etree.SubElement(pPr, qn("a:buChar"))
    bu.set("char", "•")
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")


def _clear_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    none = pPr.find(qn("a:buNone"))
    if none is None:
        etree.SubElement(pPr, qn("a:buNone"))


# ---------- shape helpers ----------

def add_rect(slide, left_in, top_in, width_in, height_in, *, fill=None,
             line=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(left_in), Inches(top_in),
                               Inches(width_in), Inches(height_in))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_width is not None:
            s.line.width = Pt(line_width)
    s.text_frame.margin_left = Emu(0)
    s.text_frame.margin_right = Emu(0)
    s.text_frame.margin_top = Emu(0)
    s.text_frame.margin_bottom = Emu(0)
    return s


def add_oval(slide, left_in, top_in, width_in, height_in, *, fill=None,
             line=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(left_in), Inches(top_in),
                               Inches(width_in), Inches(height_in))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_width is not None:
            s.line.width = Pt(line_width)
    return s


def add_line(slide, x1_in, y1_in, x2_in, y2_in, *, color, width_pt=0.75,
             dash=None):
    line = slide.shapes.add_connector(1, Inches(x1_in), Inches(y1_in),
                                      Inches(x2_in), Inches(y2_in))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash is not None:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.line.dash_style = dash
    return line


# ---------- slide chrome ----------

def add_title(slide, text, theme: Theme = DEFAULT_THEME, *, with_underline=True):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    tb = add_textbox(slide, layout.margin_left_in, layout.title_top_in,
                     width, layout.title_height_in)
    write_paragraph(tb.text_frame, text, size=typo.title_size, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    if with_underline:
        add_line(slide, layout.margin_left_in, layout.title_underline_top_in,
                 layout.slide_width_in - layout.margin_right_in,
                 layout.title_underline_top_in,
                 color=pal.rule_gray, width_pt=0.75)
    return tb


def add_subtitle_placeholder(slide, text, theme: Theme = DEFAULT_THEME, top_in=1.4):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    tb = add_textbox(slide, layout.margin_left_in, top_in, width, 0.32)
    write_paragraph(tb.text_frame, text, size=typo.body_size, italic=False,
                    color=pal.placeholder_gray, family=typo.family, first=True)
    # dashed border
    add_rect(slide, layout.margin_left_in - 0.02, top_in - 0.02,
             width + 0.04, 0.36, fill=None, line=pal.placeholder_gray,
             line_width=0.5)
    return tb


def add_section_marker(slide, label, theme: Theme = DEFAULT_THEME):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography
    w = layout.section_marker_w_in
    h = layout.section_marker_h_in
    left = layout.slide_width_in - layout.margin_right_in - w
    top = 0.18
    add_rect(slide, left, top, w, h, fill=None, line=pal.placeholder_gray,
             line_width=0.5)
    tb = add_textbox(slide, left + 0.05, top, w - 0.1, h,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, label, size=typo.small_size,
                    color=pal.footer_gray, family=typo.family,
                    align=PP_ALIGN.RIGHT, first=True)


def add_footer(slide, theme: Theme = DEFAULT_THEME, *, page_number=None,
               source: Optional[str] = None, footnote: Optional[str] = None,
               copyright_text: Optional[str] = None):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography

    foot_y = layout.footer_top_in
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    # Bottom rule line
    add_line(slide, layout.margin_left_in, foot_y,
             layout.slide_width_in - layout.margin_right_in, foot_y,
             color=pal.rule_gray, width_pt=0.5)

    # Footnote (top-left in footer)
    y = foot_y + 0.05
    if footnote:
        tb = add_textbox(slide, layout.margin_left_in, y, width / 2, 0.18)
        write_paragraph(tb.text_frame, footnote, size=typo.footer_size,
                        color=pal.text_dark, family=typo.family, first=True)
        y += 0.18
    if source:
        tb = add_textbox(slide, layout.margin_left_in, y, width / 2, 0.18)
        write_paragraph(tb.text_frame, f"Source: {source}",
                        size=typo.footer_size, color=pal.text_dark,
                        family=typo.family, first=True)

    # Right side: copyright + page number
    right_w = 4.0
    right_left = layout.slide_width_in - layout.margin_right_in - right_w
    cp = copyright_text if copyright_text is not None else theme.copyright_text
    parts = [cp]
    if page_number is not None:
        parts.append(f"   {page_number}")
    tb = add_textbox(slide, right_left, foot_y + 0.1, right_w, 0.2)
    write_paragraph(tb.text_frame, "".join(parts),
                    size=typo.footer_size, color=pal.footer_gray,
                    family=typo.family, align=PP_ALIGN.RIGHT, first=True)


def init_presentation(theme: Theme = DEFAULT_THEME):
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(theme.layout.slide_width_in)
    prs.slide_height = Inches(theme.layout.slide_height_in)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_chrome(slide, *, title, theme=DEFAULT_THEME, page_number=None,
               source=None, footnote=None, section_marker=None,
               with_underline=True):
    if section_marker:
        add_section_marker(slide, section_marker, theme)
    add_title(slide, title, theme, with_underline=with_underline)
    add_footer(slide, theme, page_number=page_number, source=source,
               footnote=footnote)
