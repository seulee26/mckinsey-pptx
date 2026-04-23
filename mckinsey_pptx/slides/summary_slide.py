"""Dark-navy full-bleed summary slide (McKinsey impact-statement style)."""
from __future__ import annotations
from typing import Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_rect, add_line, add_textbox, write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME


def add_dark_navy_summary(prs, *,
                          body: str,
                          eyebrow: Optional[str] = None,
                          page_number: Optional[int] = None,
                          corner_text: str = "McKinsey & Company",
                          theme: Theme = DEFAULT_THEME,
                          # full-bleed style — these are accepted but unused
                          section_marker=None,
                          source=None,
                          footnote=None):
    """Full-bleed dark navy slide with bold white body text.
    `body`: leading bracketed label like "[Region]: ..." is bolded by the caller
            via a leading bracket — the renderer recognises a [Foo]: prefix
            and bolds it in the same color as the rest.
    """
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Full-bleed background
    add_rect(slide, 0, 0, layout.slide_width_in, layout.slide_height_in,
             fill=pal.deep_navy)

    # Top thin rule
    add_line(slide, layout.margin_left_in, 0.7,
             layout.slide_width_in - layout.margin_right_in, 0.7,
             color=pal.white, width_pt=0.4)

    # Optional eyebrow (small, gray-ish at top right)
    if eyebrow:
        tb = add_textbox(slide, layout.margin_left_in, 0.30,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.30)
        write_paragraph(tb.text_frame, eyebrow, size=typo.small_size,
                        color=pal.placeholder_gray, family=typo.family,
                        align=PP_ALIGN.RIGHT, first=True)

    # Centered body block
    body_w = 9.5
    body_h = 4.0
    body_left = (layout.slide_width_in - body_w) / 2
    body_top = (layout.slide_height_in - body_h) / 2 - 0.2
    tb = add_textbox(slide, body_left, body_top, body_w, body_h,
                     anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Detect leading [Bracket]: prefix and bold it
    if body.startswith("[") and "]:" in body:
        end = body.index("]:") + 1
        prefix = body[:end + 1]  # includes the colon (no, end+1 is past ']')
        # Re-extract precisely: prefix is "[xxx]:"
        prefix = body[:body.index("]:") + 2]
        rest = body[len(prefix):]
        r1 = p.add_run()
        r1.text = prefix
        r1.font.size = Pt(typo.title_size + 2)
        r1.font.bold = True
        r1.font.color.rgb = pal.white
        r1.font.name = typo.family
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(typo.title_size + 2)
        r2.font.bold = True
        r2.font.color.rgb = pal.white
        r2.font.name = typo.family
    else:
        r = p.add_run()
        r.text = body
        r.font.size = Pt(typo.title_size + 2)
        r.font.bold = True
        r.font.color.rgb = pal.white
        r.font.name = typo.family

    # Corner mark (bottom right)
    cw = 3.5
    tb = add_textbox(slide, layout.slide_width_in - layout.margin_right_in - cw,
                     layout.slide_height_in - 0.4, cw, 0.25)
    parts = corner_text
    if page_number is not None:
        parts = f"{corner_text}    {page_number}"
    write_paragraph(tb.text_frame, parts, size=typo.footer_size,
                    color=pal.placeholder_gray, family=typo.family,
                    align=PP_ALIGN.RIGHT, first=True)

    return slide
