"""Timeline / roadmap / process slides:

 - phases_chevron_3: 3 chevron-arrow phases with deliverables/people lists
 - phases_chevron_4: 4 phases with timeline arrow on top + table below
 - phases_table_4: 4 phases as parallel columns with description + key activities + outcomes
 - waves_timeline_4: 4 waves on a horizontal arrow timeline with circle markers
 - gantt_timeline: workstream rows with horizontal bars by week + bottom milestones
 - overview_areas: 7 vertical area cards with letter labels
 - process_activities: 3-4 column table with Activities / Mgmt interaction / Deliverables rows
"""
from __future__ import annotations
from typing import Sequence, Optional, Dict, List, Tuple

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph, add_subtitle_placeholder,
)
from ..theme import Theme, DEFAULT_THEME


# ---------- shared chevron primitive ----------

def _chevron(slide, theme, left, top, w, h, *, fill, text=None,
             text_color=None, bold=True, size=None, align=PP_ALIGN.CENTER):
    pal, typo = theme.palette, theme.typography
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top),
                               Inches(w), Inches(h))
    s.shadow.inherit = False
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        # Use overlay textbox so we control exact alignment
        tb = add_textbox(slide, left + 0.15, top, w - 0.45, h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, text,
                        size=size or typo.body_size, bold=bold,
                        color=text_color or pal.white, family=typo.family,
                        align=align, first=True)
    return s


# ---------- 3 phases chevron ----------

def add_phases_chevron_3(prs, *,
                         title="[Project] will evolve in three major phases",
                         phases: Sequence[Dict],
                         page_number=None, section_marker="Section marker",
                         source="xx", footnote="1. xx",
                         theme: Theme = DEFAULT_THEME):
    """phases: [{"label": "Headline", "timeframe": "[Insert timeframe]",
                  "deliverables": [str, ...], "people": [str, ...]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.40

    # Top-right legend (Deliverables / People)
    leg_y = body_top
    add_oval(slide, layout.margin_left_in + width - 2.2, leg_y, 0.22, 0.22,
             fill=pal.deep_navy)
    tb = add_textbox(slide, layout.margin_left_in + width - 1.95, leg_y - 0.02,
                     1.05, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "Deliverables", size=typo.chart_label_size,
                    color=pal.text_dark, family=typo.family, first=True)
    add_oval(slide, layout.margin_left_in + width - 0.85, leg_y, 0.22, 0.22,
             fill=pal.deep_navy)
    tb = add_textbox(slide, layout.margin_left_in + width - 0.6, leg_y - 0.02,
                     0.7, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "People", size=typo.chart_label_size,
                    color=pal.text_dark, family=typo.family, first=True)

    # Timeframes row (above chevrons)
    tf_y = body_top + 0.40
    n = max(len(phases), 1)
    chev_w = (width + 0.4) / n - 0.10
    chev_h = 0.55
    chev_y = tf_y + 0.30
    chev_left0 = layout.margin_left_in
    chev_xs = [chev_left0 + i * (chev_w - 0.40) for i in range(n)]

    for x, ph in zip(chev_xs, phases):
        tb = add_textbox(slide, x + 0.20, tf_y, chev_w - 0.20, 0.28)
        write_paragraph(tb.text_frame,
                        f"[{ph.get('timeframe', 'Insert timeframe')}]",
                        size=typo.chart_label_size,
                        color=pal.placeholder_gray,
                        family=typo.family, first=True)

    # Chevrons (alternating bright + deep navy as in source)
    chev_colors = [pal.bright_blue, pal.deep_navy, pal.deep_navy,
                   pal.deep_navy, pal.deep_navy]
    for i, (x, ph) in enumerate(zip(chev_xs, phases)):
        fill = chev_colors[min(i, len(chev_colors) - 1)]
        # number circle inside chevron, then label after it
        _chevron(slide, theme, x, chev_y, chev_w, chev_h, fill=fill,
                 text=None)
        # Number circle
        n_d = 0.36
        nx = x + 0.18
        ny = chev_y + (chev_h - n_d) / 2
        add_oval(slide, nx, ny, n_d, n_d, fill=pal.white)
        tb = add_textbox(slide, nx, ny, n_d, n_d, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(i + 1), size=typo.body_size,
                        bold=True, color=pal.deep_navy, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Label
        tb = add_textbox(slide, nx + n_d + 0.10, chev_y,
                         chev_w - n_d - 0.55, chev_h, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, ph.get("label", "[Phase headline]"),
                        size=typo.body_size, bold=True, color=pal.white,
                        family=typo.family, first=True)

    # Bottom area split: deliverables (top) + people (bottom)
    body_bottom = layout.footer_top_in - 0.20
    body_h = body_bottom - (chev_y + chev_h + 0.2)
    deliv_top = chev_y + chev_h + 0.30
    people_top = deliv_top + body_h * 0.55
    for i, (x, ph) in enumerate(zip(chev_xs, phases)):
        col_left = x + 0.20
        col_w = chev_w - 0.40
        # Deliverables icon
        add_oval(slide, col_left, deliv_top, 0.20, 0.20, fill=pal.deep_navy)
        tb = add_textbox(slide, col_left + 0.30, deliv_top - 0.04,
                         col_w - 0.30, body_h * 0.50)
        first = True
        for d in ph.get("deliverables", []):
            write_paragraph(tb.text_frame, d, size=typo.body_size - 1,
                            color=pal.text_dark, family=typo.family,
                            space_after=2, first=first)
            first = False
        # People icon
        add_oval(slide, col_left, people_top, 0.20, 0.20, fill=pal.deep_navy)
        tb = add_textbox(slide, col_left + 0.30, people_top - 0.04,
                         col_w - 0.30, body_bottom - people_top)
        first = True
        for pp in ph.get("people", []):
            write_paragraph(tb.text_frame, pp, size=typo.body_size - 1,
                            color=pal.text_dark, family=typo.family,
                            space_after=2, first=first)
            first = False
    return slide


# ---------- 4 phases as columns (text-heavy) ----------

def add_phases_table_4(prs, *,
                       title="[Project] will evolve through four main phases",
                       subtitle: Optional[str] = "[Insert subtitle]",
                       phases: Sequence[Dict],
                       page_number=None, section_marker="Section marker",
                       source="xx", footnote="1. xx",
                       theme: Theme = DEFAULT_THEME):
    """phases: [{"name": "Phase 1", "description": "...",
                 "activities": [...], "outcomes": [...]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.95
    body_bottom = layout.footer_top_in - 0.20

    n = max(len(phases), 1)
    col_w = (width - (n - 1) * 0.25) / n
    arrow_y = body_top + 0.10

    for i, ph in enumerate(phases):
        cx = layout.margin_left_in + i * (col_w + 0.25)
        # [TIMELINE] marker
        tb = add_textbox(slide, cx, arrow_y, col_w, 0.20)
        write_paragraph(tb.text_frame, "[TIMELINE]",
                        size=typo.chart_label_size,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)
        # Phase name (bright blue)
        tb = add_textbox(slide, cx, arrow_y + 0.25, col_w, 0.40)
        write_paragraph(tb.text_frame, ph.get("name", f"Phase {i+1}"),
                        size=typo.title_size - 4, bold=True,
                        color=pal.bright_blue, family=typo.family, first=True)
        # Right arrow indicator (small)
        if i < n - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         Inches(cx + col_w + 0.02),
                                         Inches(arrow_y + 0.40),
                                         Inches(0.18), Inches(0.18))
            arr.shadow.inherit = False
            arr.fill.solid(); arr.fill.fore_color.rgb = pal.placeholder_gray
            arr.line.fill.background()

        # Description
        tb = add_textbox(slide, cx, arrow_y + 0.70, col_w, 1.2)
        write_paragraph(tb.text_frame,
                        ph.get("description",
                               "[High-level description of main purpose of phase]"),
                        size=typo.body_size, color=pal.text_dark,
                        family=typo.family, first=True)

        # Key activities header + bullets
        sec_top = arrow_y + 1.95
        tb = add_textbox(slide, cx, sec_top, col_w, 0.30)
        write_paragraph(tb.text_frame, "Key activities",
                        size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        tb = add_textbox(slide, cx, sec_top + 0.30, col_w, 1.6)
        first = True
        for a in ph.get("activities", []):
            write_paragraph(tb.text_frame, a, size=typo.body_size - 1,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
        # Outcomes header + bullets
        out_top = sec_top + 1.95
        tb = add_textbox(slide, cx, out_top, col_w, 0.30)
        write_paragraph(tb.text_frame, "Outcomes",
                        size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        tb = add_textbox(slide, cx, out_top + 0.30, col_w,
                         body_bottom - out_top - 0.30)
        first = True
        for o in ph.get("outcomes", []):
            write_paragraph(tb.text_frame, o, size=typo.body_size - 1,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
    return slide


# ---------- 4 waves on horizontal timeline ----------

def add_waves_timeline_4(prs, *,
                         title="[Project] will be built in four main waves",
                         subtitle: Optional[str] = "[Insert subtitle]",
                         waves: Sequence[Dict],
                         page_number=None, section_marker="Section marker",
                         source="xx", footnote="1. xx",
                         theme: Theme = DEFAULT_THEME):
    """waves: [{"name": "Wave 1", "headline": "[Insert headline]",
                 "timeframe": "[Insert timeframe]",
                 "activities": [...], "deliverables": [...]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 2.10
    n = max(len(waves), 1)

    # Horizontal arrow with markers
    arrow_y = body_top + 0.85
    add_line(slide, layout.margin_left_in, arrow_y,
             layout.margin_left_in + width - 0.5, arrow_y,
             color=pal.text_dark, width_pt=1.0)
    # Arrow head (right end)
    head = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                  Inches(layout.margin_left_in + width - 0.55),
                                  Inches(arrow_y - 0.10),
                                  Inches(0.18), Inches(0.20))
    head.shadow.inherit = False
    head.fill.solid(); head.fill.fore_color.rgb = pal.text_dark
    head.line.fill.background()

    slot = (width - 0.6) / n
    for i, w in enumerate(waves):
        cx = layout.margin_left_in + i * slot + slot / 2 - 0.4
        # Timeframe (top)
        tb = add_textbox(slide, cx - slot / 2 + 0.3, body_top,
                         slot - 0.3, 0.25)
        write_paragraph(tb.text_frame,
                        f"[{w.get('timeframe', 'Insert timeframe')}]",
                        size=typo.chart_label_size,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)
        # Wave name + headline
        tb = add_textbox(slide, cx - slot / 2 + 0.3, body_top + 0.30,
                         slot - 0.3, 0.30)
        write_paragraph(tb.text_frame, w.get("name", f"Wave {i+1}"),
                        size=typo.body_size, bold=True,
                        color=pal.bright_blue, family=typo.family, first=True)
        tb = add_textbox(slide, cx - slot / 2 + 0.3, body_top + 0.55,
                         slot - 0.3, 0.30)
        write_paragraph(tb.text_frame,
                        w.get("headline", "[Insert headline]"),
                        size=typo.body_size - 1, color=pal.text_dark,
                        family=typo.family, first=True)
        # Marker on arrow
        md = 0.30
        marker_cx = layout.margin_left_in + i * slot + slot / 2
        add_oval(slide, marker_cx - md / 2, arrow_y - md / 2, md, md,
                 fill=pal.bright_blue)

        # Body sections
        col_left = layout.margin_left_in + i * slot + 0.10
        col_w = slot - 0.20

        # Key activities
        ka_top = arrow_y + 0.35
        tb = add_textbox(slide, col_left, ka_top, col_w, 0.30)
        write_paragraph(tb.text_frame, "Key activities:",
                        size=typo.body_size - 1, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        tb = add_textbox(slide, col_left, ka_top + 0.30, col_w, 1.5)
        first = True
        for a in w.get("activities", []):
            write_paragraph(tb.text_frame, a, size=typo.body_size - 2,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False

        md_top = ka_top + 1.85
        tb = add_textbox(slide, col_left, md_top, col_w, 0.30)
        write_paragraph(tb.text_frame, "Main deliverables:",
                        size=typo.body_size - 1, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        tb = add_textbox(slide, col_left, md_top + 0.30, col_w,
                         layout.footer_top_in - md_top - 0.50)
        first = True
        for d in w.get("deliverables", []):
            write_paragraph(tb.text_frame, d, size=typo.body_size - 2,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
    return slide


# ---------- Gantt-style timeline ----------

def add_gantt_timeline(prs, *,
                       title="High-level timeline for next [insert timeframe]",
                       subtitle: Optional[str] = "[Insert subtitle]",
                       weeks: Sequence,
                       workstreams: Sequence[Dict],
                       milestones: Sequence[Dict] = (),
                       page_number=None, section_marker="Section marker",
                       source="xx", footnote="1. xx",
                       theme: Theme = DEFAULT_THEME):
    """workstreams: [{"name": "01 [Insert work stream]", "start_week": int,
                       "end_week": int, "color": "blue_light"|"blue_dark"}]
       milestones: [{"week": int, "label": "Kick-off"}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.95

    # Layout: left labels column, right timeline grid
    label_w = 1.7
    grid_left = layout.margin_left_in + label_w
    grid_right = layout.slide_width_in - layout.margin_right_in
    grid_w = grid_right - grid_left

    # Top header row
    header_y = body_top
    n_weeks = len(weeks)
    if n_weeks == 0:
        return slide
    cell_w = grid_w / n_weeks
    # "Week" label
    tb = add_textbox(slide, layout.margin_left_in, header_y, label_w - 0.1, 0.28)
    write_paragraph(tb.text_frame, "Week", size=typo.body_size - 1, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    for i, wk in enumerate(weeks):
        tb = add_textbox(slide, grid_left + i * cell_w,
                         header_y, cell_w, 0.28,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(wk),
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    # Workstream rows
    rows_top = header_y + 0.40
    rows_bottom = layout.footer_top_in - 1.10  # leave space for milestone labels
    n_rows = max(len(workstreams), 1)
    row_h = (rows_bottom - rows_top) / n_rows

    color_map = {
        "blue_light": pal.bright_blue,
        "blue_dark": pal.deep_navy,
        "navy": pal.dark_navy,
        "royal": pal.royal_blue,
    }

    for ri, ws in enumerate(workstreams):
        ry = rows_top + ri * row_h
        # Label
        tb = add_textbox(slide, layout.margin_left_in, ry, label_w - 0.1, row_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, ws.get("name", "[Insert work stream]"),
                        size=typo.chart_label_size, color=pal.text_dark,
                        family=typo.family, first=True)
        # Vertical week separators (light)
        for i in range(n_weeks + 1):
            add_line(slide, grid_left + i * cell_w, ry,
                     grid_left + i * cell_w, ry + row_h,
                     color=pal.grid_gray, width_pt=0.3)
        # Bar
        sw = ws.get("start_week")
        ew = ws.get("end_week")
        if sw is None or ew is None:
            continue
        try:
            si = weeks.index(sw)
            ei = weeks.index(ew)
        except ValueError:
            continue
        bar_left = grid_left + si * cell_w + 0.04
        bar_right = grid_left + (ei + 1) * cell_w - 0.04
        bar_h = row_h * 0.55
        bar_y = ry + (row_h - bar_h) / 2
        fill = color_map.get(ws.get("color", "blue_light"), pal.bright_blue)
        add_rect(slide, bar_left, bar_y, bar_right - bar_left, bar_h,
                 fill=fill)

    # Milestones row (bottom)
    ms_y = rows_bottom + 0.05
    label_y = ms_y + 0.30
    md = 0.22
    for ms in milestones:
        wk = ms.get("week")
        if wk is None:
            continue
        try:
            wi = weeks.index(wk)
        except ValueError:
            continue
        cx = grid_left + wi * cell_w + cell_w / 2
        add_oval(slide, cx - md / 2, ms_y, md, md, fill=pal.status_amber)
        tb = add_textbox(slide, cx - 1.0, label_y, 2.0, 0.45)
        write_paragraph(tb.text_frame, ms.get("label", ""),
                        size=typo.chart_label_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
    return slide


# ---------- Overview of N areas (vertical cards with letter badge) ----------

def add_overview_areas(prs, *,
                       title="[Overview of areas / Insert action title]",
                       subtitle: Optional[str] = "[Insert subtitle]",
                       areas: Sequence[Dict],
                       call_out: Optional[str] = None,
                       page_number=None, section_marker="Section marker",
                       source="xx", footnote="1. xx",
                       theme: Theme = DEFAULT_THEME):
    """areas: [{"name": "[Area 1]", "bullets": [str, ...]}]
    Up to ~7 cards.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 2.05
    body_bottom = layout.footer_top_in - 0.20

    n = max(len(areas), 1)
    gap = 0.12
    card_w = (width - (n - 1) * gap) / n

    badge_d = 0.42
    head_h = 0.50
    card_h = body_bottom - body_top - badge_d / 2

    letters = "ABCDEFGHIJ"

    for i, area in enumerate(areas):
        cx = layout.margin_left_in + i * (card_w + gap)
        card_top = body_top + badge_d / 2
        # Card body (light gray)
        add_rect(slide, cx, card_top, card_w, card_h, fill=pal.soft_gray)
        # Header band (deep navy)
        add_rect(slide, cx, card_top, card_w, head_h, fill=pal.deep_navy)
        tb = add_textbox(slide, cx + 0.10, card_top, card_w - 0.20, head_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, area.get("name", f"[Area {i+1}]"),
                        size=typo.body_size, bold=True, color=pal.white,
                        family=typo.family, first=True)
        # Letter badge (circle straddling top edge)
        bx = cx + (card_w - badge_d) / 2
        by = body_top
        add_oval(slide, bx, by, badge_d, badge_d, fill=pal.white,
                 line=pal.deep_navy, line_width=1.0)
        tb = add_textbox(slide, bx, by, badge_d, badge_d,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, letters[i] if i < len(letters)
                        else str(i + 1),
                        size=typo.body_size, bold=True, color=pal.deep_navy,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
        # Bullets
        tb = add_textbox(slide, cx + 0.10, card_top + head_h + 0.05,
                         card_w - 0.20, card_h - head_h - 0.10)
        first = True
        for b in area.get("bullets", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size - 2,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False

    # Optional call-out (bottom-left blue tag)
    if call_out:
        co_w = 2.4
        co_h = 0.7
        co_x = layout.margin_left_in
        co_y = body_bottom - co_h - 0.15
        add_rect(slide, co_x, co_y, co_w, co_h, fill=pal.bright_blue)
        tb = add_textbox(slide, co_x + 0.10, co_y, co_w - 0.20, co_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, call_out,
                        size=typo.body_size - 1, color=pal.white,
                        family=typo.family, first=True)
    return slide


# ---------- Process activities table ----------

def add_process_activities(prs, *,
                           title="[Project title] process expected to run [xx weeks/months]",
                           subtitle: Optional[str] = "High-level project plan and deliverables",
                           steps: Sequence[Dict],
                           page_number=None, section_marker="Section marker",
                           source="xx", footnote="1. xx",
                           theme: Theme = DEFAULT_THEME):
    """steps: [{"name": "Week xx", "subtitle": "Research and hypotheses",
                "activities": [str, ...],
                "interaction": str,           # short label like "Kick-off"
                "deliverable": str | None}]   # short label or None
    Renders 3-row layout: Activities / Mgmt. interaction / Deliverables
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        # Use a regular textbox (not the dashed placeholder)
        tb = add_textbox(slide, layout.margin_left_in, 1.30,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.30)
        write_paragraph(tb.text_frame, subtitle, size=typo.section_title_size,
                        bold=True, color=pal.text_dark,
                        family=typo.family, first=True)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.75

    # Left fixed gutter for row labels
    label_w = 1.4
    grid_left = layout.margin_left_in + label_w
    grid_w = width - label_w
    n = max(len(steps), 1)
    col_w = grid_w / n

    # Header row for each step (deep-navy band)
    head_h = 0.55
    for i, st in enumerate(steps):
        cx = grid_left + i * col_w
        add_rect(slide, cx + 0.05, body_top, col_w - 0.10, head_h,
                 fill=pal.deep_navy)
        tb = add_textbox(slide, cx + 0.10, body_top + 0.03, col_w - 0.20, 0.24,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, st.get("name", f"Week {i+1}"),
                        size=typo.body_size - 1, bold=True, color=pal.white,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
        tb = add_textbox(slide, cx + 0.10, body_top + 0.27, col_w - 0.20, 0.24,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, st.get("subtitle", ""),
                        size=typo.chart_label_size, color=pal.white,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    # Row labels
    rows = [
        ("Activities", 2.5),
        ("Mgmt. interaction", 0.7),
        ("Deliverables", 0.7),
    ]
    ry = body_top + head_h + 0.10
    for lbl, h in rows:
        tb = add_textbox(slide, layout.margin_left_in, ry, label_w - 0.1, 0.30)
        write_paragraph(tb.text_frame, lbl, size=typo.body_size - 1,
                        bold=True, color=pal.text_dark, family=typo.family,
                        first=True)
        # Light bottom rule
        add_line(slide, layout.margin_left_in,
                 ry + h, layout.slide_width_in - layout.margin_right_in,
                 ry + h, color=pal.grid_gray, width_pt=0.4)
        ry += h

    # Activities row content
    a_top = body_top + head_h + 0.10
    for i, st in enumerate(steps):
        cx = grid_left + i * col_w + 0.10
        tb = add_textbox(slide, cx, a_top, col_w - 0.20, 2.4)
        first = True
        for a in st.get("activities", []):
            write_paragraph(tb.text_frame, a, size=typo.body_size - 2,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False

    # Interaction row content (diamond marker + label)
    int_top = a_top + 2.5
    for i, st in enumerate(steps):
        cx = grid_left + i * col_w + col_w / 2
        if st.get("interaction"):
            d = 0.22
            shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                           Inches(cx - d / 2),
                                           Inches(int_top + 0.05),
                                           Inches(d), Inches(d))
            shape.shadow.inherit = False
            shape.fill.solid(); shape.fill.fore_color.rgb = pal.deep_navy
            shape.line.fill.background()
            tb = add_textbox(slide, cx + d / 2 + 0.05, int_top + 0.04,
                             col_w / 2 + 0.5, 0.28,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, st["interaction"],
                            size=typo.body_size - 1, color=pal.text_dark,
                            family=typo.family, first=True)

    # Deliverables row content
    del_top = int_top + 0.7
    for i, st in enumerate(steps):
        cx = grid_left + i * col_w + col_w / 2
        if st.get("deliverable"):
            d = 0.22
            shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                           Inches(cx - d / 2),
                                           Inches(del_top + 0.05),
                                           Inches(d), Inches(d))
            shape.shadow.inherit = False
            shape.fill.solid(); shape.fill.fore_color.rgb = pal.bright_blue
            shape.line.fill.background()
            tb = add_textbox(slide, cx + d / 2 + 0.05, del_top + 0.04,
                             col_w / 2 + 0.5, 0.28,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, st["deliverable"],
                            size=typo.body_size - 1, color=pal.text_dark,
                            family=typo.family, first=True)
    return slide
