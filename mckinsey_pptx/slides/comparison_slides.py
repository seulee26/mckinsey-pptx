"""Comparison-style slides:
 - comparison_table: criteria rows × option columns, with Harvey-ball ratings
 - pros_cons: 2-column layout (✓ green / ✗ red)
 - two_column_compare: Before/After or As-is/To-be side-by-side
"""
from __future__ import annotations
from typing import Sequence, Optional, Literal
from math import pi

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME


# Harvey-ball levels: 0 = empty, 4 = full
def _harvey_ball(slide, theme, cx, cy, d, level: int):
    """Draw a 0-4 quartile-fill Harvey ball centered at (cx, cy)."""
    pal = theme.palette
    level = max(0, min(int(level), 4))
    # Outer outline
    add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=pal.white,
             line=pal.deep_navy, line_width=1.0)
    if level == 0:
        return
    if level == 4:
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=pal.deep_navy,
                 line=pal.deep_navy, line_width=1.0)
        return
    # Half / quarter / three-quarter via overlapping pies (approximate w/ rects + ovals)
    # Strategy: draw filled circle clipped by white rectangle(s)
    # Half (level 2): cover right half with white rect
    # Quarter (level 1): cover bottom-right (270° kept empty + top-left)
    # 3/4 (level 3): cover top-right quarter only
    if level == 2:
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=pal.deep_navy,
                 line=pal.deep_navy, line_width=1.0)
        add_rect(slide, cx, cy - d / 2 - 0.005, d / 2 + 0.005, d + 0.010,
                 fill=pal.white)
        # Re-stroke outline with thin oval line
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=None,
                 line=pal.deep_navy, line_width=1.0)
        return
    if level == 1:
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=pal.deep_navy,
                 line=pal.deep_navy, line_width=1.0)
        # Cover right half + bottom-left quarter -> only top-left filled
        add_rect(slide, cx, cy - d / 2 - 0.005, d / 2 + 0.005, d + 0.010,
                 fill=pal.white)
        add_rect(slide, cx - d / 2 - 0.005, cy, d / 2 + 0.005, d / 2 + 0.010,
                 fill=pal.white)
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=None,
                 line=pal.deep_navy, line_width=1.0)
        return
    if level == 3:
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=pal.deep_navy,
                 line=pal.deep_navy, line_width=1.0)
        # Erase top-right quarter only
        add_rect(slide, cx, cy - d / 2 - 0.005, d / 2 + 0.005, d / 2 + 0.005,
                 fill=pal.white)
        add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=None,
                 line=pal.deep_navy, line_width=1.0)
        return


def _parse_score(score):
    """Accept int 0-4, str like '●●●○○' (count filled balls), or
    semantic 'high'|'med'|'low'|'none'."""
    if isinstance(score, int):
        return max(0, min(score, 4))
    if isinstance(score, str):
        s = score.strip().lower()
        if s in ("high", "strong", "best", "yes"):
            return 4
        if s in ("med-high", "good"):
            return 3
        if s in ("medium", "med", "ok", "moderate"):
            return 2
        if s in ("low", "weak", "limited"):
            return 1
        if s in ("none", "no", "n/a", "-"):
            return 0
        # Otherwise count filled circles in the string
        if any(ch in s for ch in "●◐○"):
            full = s.count("●")
            half = s.count("◐")
            return min(4, int(full * 1 + half * 0.5))
    return 0


# ---------- Comparison table ----------

def add_comparison_table(prs, *,
                         title="[Option comparison / Insert action title]",
                         subtitle: Optional[str] = None,
                         options: Sequence[str],
                         criteria: Sequence[dict],
                         recommended_index: Optional[int] = None,
                         page_number=None, section_marker="Section marker",
                         source="xx", footnote="1. xx",
                         theme: Theme = DEFAULT_THEME):
    """Side-by-side option comparison.
       options: list of column headers (e.g. ["Option A", "Option B", "Option C"])
       criteria: list of rows, each:
                 {"name": "Cost", "scores": [4, 2, 3], "notes": ["Low", "High", "Med"]}
                 - scores: int 0-4 OR semantic str OR Harvey-ball glyphs
                 - notes (optional): per-option short text shown under the ball
       recommended_index: optional column to highlight (light blue header band)
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.50
    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.30)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size, color=pal.placeholder_gray,
                        family=typo.family, first=True)
        body_top += 0.40

    body_bottom = layout.footer_top_in - 0.20
    n_opts = max(len(options), 1)
    crit_w = 2.4
    opt_w = (width - crit_w) / n_opts

    # Header row
    head_h = 0.55
    add_rect(slide, layout.margin_left_in, body_top, crit_w, head_h,
             fill=pal.soft_gray)
    tb = add_textbox(slide, layout.margin_left_in + 0.10, body_top,
                     crit_w - 0.20, head_h, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "Criteria", size=typo.body_size, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    for i, opt in enumerate(options):
        ox = layout.margin_left_in + crit_w + i * opt_w
        fill = pal.bright_blue if recommended_index == i else pal.deep_navy
        add_rect(slide, ox, body_top, opt_w, head_h, fill=fill)
        tb = add_textbox(slide, ox + 0.10, body_top, opt_w - 0.20, head_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, opt, size=typo.body_size, bold=True,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

    # Body rows
    rows_top = body_top + head_h
    n_rows = max(len(criteria), 1)
    row_h = (body_bottom - rows_top - 0.10) / n_rows

    ball_d = min(0.45, row_h * 0.55)

    for ri, row in enumerate(criteria):
        ry = rows_top + ri * row_h
        # alternating row tint
        if ri % 2 == 0:
            add_rect(slide,
                     layout.margin_left_in + crit_w, ry,
                     opt_w * n_opts, row_h, fill=pal.soft_gray)
        # Criteria column
        tb = add_textbox(slide, layout.margin_left_in + 0.10, ry,
                         crit_w - 0.20, row_h, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, row.get("name", ""),
                        size=typo.body_size, bold=True, color=pal.text_dark,
                        family=typo.family, first=True)
        scores = row.get("scores", [])
        notes = row.get("notes", [])
        for i in range(n_opts):
            ox = layout.margin_left_in + crit_w + i * opt_w
            cx = ox + opt_w / 2
            note = notes[i] if i < len(notes) else None
            score = scores[i] if i < len(scores) else None
            if score is not None:
                _harvey_ball(slide, theme, cx,
                             ry + (row_h - (0.30 if note else 0)) / 2,
                             ball_d, _parse_score(score))
            if note:
                tb = add_textbox(slide, ox + 0.10, ry + row_h - 0.36,
                                 opt_w - 0.20, 0.30,
                                 anchor=MSO_ANCHOR.MIDDLE)
                write_paragraph(tb.text_frame, note,
                                size=typo.body_size - 2,
                                color=pal.text_dark, family=typo.family,
                                align=PP_ALIGN.CENTER, first=True)
        # Bottom rule
        add_line(slide, layout.margin_left_in, ry + row_h,
                 layout.margin_left_in + width, ry + row_h,
                 color=pal.grid_gray, width_pt=0.4)

    # Optional "Recommended" tag below highlighted column
    if recommended_index is not None:
        ox = layout.margin_left_in + crit_w + recommended_index * opt_w
        tb = add_textbox(slide, ox, body_bottom - 0.05, opt_w, 0.25)
        write_paragraph(tb.text_frame, "★ Recommended",
                        size=typo.chart_label_size, bold=True,
                        color=pal.bright_blue, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
    return slide


# ---------- Pros & Cons ----------

def add_pros_cons(prs, *,
                  title="[Pros and cons / Insert action title]",
                  subtitle: Optional[str] = None,
                  pros: Sequence[str],
                  cons: Sequence[str],
                  pros_label: str = "Pros",
                  cons_label: str = "Cons",
                  page_number=None, section_marker="Section marker",
                  source="xx", footnote="1. xx",
                  theme: Theme = DEFAULT_THEME):
    """Two-column pros/cons. Green ✓ on left, red ✗ on right."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.50
    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.30)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size, color=pal.placeholder_gray,
                        family=typo.family, first=True)
        body_top += 0.45

    body_bottom = layout.footer_top_in - 0.20
    body_h = body_bottom - body_top
    col_w = (width - 0.40) / 2

    sides = [
        (layout.margin_left_in, pros_label, pros, pal.status_green, "✓"),
        (layout.margin_left_in + col_w + 0.40, cons_label, cons,
         pal.status_red, "✗"),
    ]
    for left, label, items, color, glyph in sides:
        # Header band
        head_h = 0.55
        add_rect(slide, left, body_top, col_w, head_h, fill=color)
        tb = add_textbox(slide, left + 0.20, body_top, col_w - 0.20, head_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, label, size=typo.section_title_size + 2,
                        bold=True, color=pal.white, family=typo.family,
                        first=True)
        # Body card (light gray)
        card_top = body_top + head_h
        card_h = body_h - head_h
        add_rect(slide, left, card_top, col_w, card_h, fill=pal.soft_gray)
        # Items
        ny = card_top + 0.20
        item_h = max(0.45, (card_h - 0.40) / max(len(items), 1))
        for it in items:
            # glyph
            tb = add_textbox(slide, left + 0.20, ny, 0.35, item_h,
                             anchor=MSO_ANCHOR.TOP)
            write_paragraph(tb.text_frame, glyph,
                            size=typo.body_size + 2, bold=True,
                            color=color, family=typo.family, first=True)
            # text
            tb = add_textbox(slide, left + 0.60, ny,
                             col_w - 0.80, item_h, anchor=MSO_ANCHOR.TOP)
            write_paragraph(tb.text_frame, it, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            first=True)
            ny += item_h
    return slide


# ---------- Two-column compare ----------

def add_two_column_compare(prs, *,
                           title="[Before / After / Insert action title]",
                           subtitle: Optional[str] = None,
                           left_label: str = "Before",
                           right_label: str = "After",
                           left_items: Sequence[str],
                           right_items: Sequence[str],
                           left_color: Literal["navy", "gray", "blue", "amber"] = "gray",
                           right_color: Literal["navy", "gray", "blue", "amber"] = "blue",
                           show_arrow: bool = True,
                           page_number=None, section_marker="Section marker",
                           source="xx", footnote="1. xx",
                           theme: Theme = DEFAULT_THEME):
    """Before/After or As-is/To-be comparison. Two cards with a connecting arrow."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    color_map = {
        "navy": pal.deep_navy,
        "gray": pal.footer_gray,
        "blue": pal.bright_blue,
        "amber": pal.status_amber,
    }
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.50
    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.30)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size, color=pal.placeholder_gray,
                        family=typo.family, first=True)
        body_top += 0.45

    body_bottom = layout.footer_top_in - 0.20
    body_h = body_bottom - body_top
    arrow_w = 0.9 if show_arrow else 0.4
    col_w = (width - arrow_w) / 2

    sides = [
        (layout.margin_left_in, left_label, left_items, color_map[left_color]),
        (layout.margin_left_in + col_w + arrow_w, right_label, right_items,
         color_map[right_color]),
    ]
    for left, label, items, color in sides:
        head_h = 0.55
        add_rect(slide, left, body_top, col_w, head_h, fill=color)
        tb = add_textbox(slide, left + 0.20, body_top, col_w - 0.20, head_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, label, size=typo.section_title_size + 2,
                        bold=True, color=pal.white, family=typo.family,
                        first=True)
        card_top = body_top + head_h
        card_h = body_h - head_h
        add_rect(slide, left, card_top, col_w, card_h, fill=pal.soft_gray)
        # Bullets
        tb = add_textbox(slide, left + 0.20, card_top + 0.15,
                         col_w - 0.40, card_h - 0.30)
        first = True
        for it in items:
            write_paragraph(tb.text_frame, it, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=4, first=first)
            first = False

    if show_arrow:
        # Big arrow between the two columns
        arr_x = layout.margin_left_in + col_w + 0.10
        arr_y = body_top + body_h / 2 - 0.40
        arr = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(arr_x), Inches(arr_y),
            Inches(arrow_w - 0.20), Inches(0.80))
        arr.shadow.inherit = False
        arr.fill.solid(); arr.fill.fore_color.rgb = pal.bright_blue
        arr.line.fill.background()
    return slide
