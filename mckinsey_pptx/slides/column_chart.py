"""Column chart slides — 4 variants matching the McKinsey templates.

Layout: left ~ chart with [Description] header, right column = takeaways.
Charts are drawn with native shapes (not pptx chart) so styling matches exactly.
"""
from __future__ import annotations
from typing import Sequence, Optional, Tuple
import math

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME


# Geometry constants (relative to slide left/top in inches)
DEFAULT_CHART_BOX = (0.45, 1.95, 8.5, 4.85)  # left, top, width, height
DEFAULT_TAKEAWAY_BOX = (9.45, 1.95, 3.45, 4.85)
TAKEAWAY_DIVIDER_X = 9.30
DEFAULT_DESCRIPTION_TOP = 1.45  # below title underline (1.15) with breathing room


def _y_ticks(max_val, target_ticks=9):
    """Produce a clean (1/2/2.5/5)*10^k step that gives roughly target_ticks ticks
    spanning [0, top] where top >= max_val * 1.05."""
    if max_val <= 0:
        return [0, 1], 1
    target = max_val * 1.05
    raw_step = target / (target_ticks - 1)
    mag = 10 ** math.floor(math.log10(raw_step))
    best = None
    for mult in (1, 2, 2.5, 5, 10, 20):
        step = mult * mag
        n = math.ceil(target / step) + 1
        if 5 <= n <= 12:
            score = abs(n - target_ticks)
            if best is None or score < best[0]:
                best = (score, step, n)
    if best is None:
        step = mag * 10
        n = math.ceil(target / step) + 1
        best = (0, step, n)
    _, step, n = best
    return [step * i for i in range(n)], step * (n - 1)


def _draw_takeaway(slide, theme, *, top=None, takeaways: Sequence[str],
                   header="Key takeaways/main conclusion",
                   box=DEFAULT_TAKEAWAY_BOX,
                   divider_x=TAKEAWAY_DIVIDER_X):
    pal, typo = theme.palette, theme.typography
    left, ttop, w, h = box
    if top is not None:
        ttop = top
    # Dashed vertical divider
    add_line(slide, divider_x, ttop, divider_x, ttop + h,
             color=pal.placeholder_gray, width_pt=0.5)
    # Header — placeholder default = gray `[...]`; real content = dark text.
    if header == "Key takeaways/main conclusion":
        header_text = f"[{header}]"
        header_color = pal.placeholder_gray
    elif header.startswith("["):
        header_text = header
        header_color = pal.placeholder_gray
    else:
        header_text = header
        header_color = pal.text_dark
    tb = add_textbox(slide, left, ttop - 0.05, w, 0.35)
    write_paragraph(tb.text_frame, header_text, size=typo.section_title_size,
                    bold=True, color=header_color, family=typo.family,
                    first=True)
    # Header underline
    add_line(slide, left, ttop + 0.30, left + w, ttop + 0.30,
             color=pal.rule_gray, width_pt=0.5)
    # Bullets
    tb = add_textbox(slide, left, ttop + 0.55, w, h - 0.55)
    first = True
    for t in takeaways:
        write_paragraph(tb.text_frame, t, size=typo.body_size,
                        color=pal.text_dark, family=typo.family,
                        bullet=True, space_after=10, first=first)
        first = False


def _draw_description_header(slide, theme, *, left, top, width, label="Description"):
    """Render the '[Description]' header above a chart.

    - `label=None` or empty string: render only the rule, no text.
    - `label` starts with `[`: render as-is (caller controls bracketing).
    - `label == "Description"` (template default): render `[Description]` in
      placeholder styling (gray) to signal the slot is unfilled.
    - Any other string: render as a real, dark-text header.
    """
    pal, typo = theme.palette, theme.typography
    if label is None or label == "":
        # Just the rule, no header text.
        add_line(slide, left, top + 0.30, left + width, top + 0.30,
                 color=pal.rule_gray, width_pt=0.5)
        return
    if label == "Description":
        label_text = f"[{label}]"
        color = pal.placeholder_gray
    elif label.startswith("["):
        label_text = label
        color = pal.placeholder_gray
    else:
        label_text = label
        color = pal.text_dark
    tb = add_textbox(slide, left, top - 0.05, width, 0.35)
    write_paragraph(tb.text_frame, label_text,
                    size=typo.section_title_size, bold=True,
                    color=color, family=typo.family, first=True)
    add_line(slide, left, top + 0.30, left + width, top + 0.30,
             color=pal.rule_gray, width_pt=0.5)


def _draw_axis_and_bars(slide, theme, *, chart_box, data_label,
                        data_unit, categories, values, focus_index=None,
                        forecast_from_index=None, legend=None):
    """Returns (x_for_each_bar_center, baseline_y, bar_w, top_y, max_val).

    Negative values not supported (templates don't show them)."""
    pal, typo = theme.palette, theme.typography
    cleft, ctop, cwidth, cheight = chart_box

    # Title above chart: "Data, Unit" — bracket+gray for template defaults,
    # dark+unbracketed for real content.
    lbl_ph = data_label == "Data"
    unit_ph = data_unit == "Unit"
    lbl_text = f"[{data_label}]" if lbl_ph else data_label
    unit_text = f"[{data_unit}]" if unit_ph else data_unit
    tb = add_textbox(slide, cleft, ctop, cwidth - 1.5, 0.3)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = f"{lbl_text}, "
    r.font.size = Pt(typo.section_title_size); r.font.bold = True
    r.font.color.rgb = pal.text_dark; r.font.name = typo.family
    r2 = p.add_run(); r2.text = unit_text
    r2.font.size = Pt(typo.section_title_size)
    r2.font.color.rgb = pal.placeholder_gray; r2.font.name = typo.family

    # Legend (right side of chart top)
    if legend:
        leg_y = ctop
        leg_left = cleft + cwidth - 0.05
        for color_rgb, label in reversed(legend):
            label_w = 0.7
            swatch_w = 0.22
            tb = add_textbox(slide, leg_left - label_w, leg_y, label_w, 0.25,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, label, size=typo.chart_label_size,
                            color=pal.text_dark, family=typo.family,
                            align=PP_ALIGN.LEFT, first=True)
            add_rect(slide, leg_left - label_w - swatch_w - 0.05,
                     leg_y + 0.05, swatch_w, 0.15, fill=color_rgb)
            leg_left -= label_w + swatch_w + 0.25

    # Compute axis area
    plot_left = cleft + 0.55
    plot_top = ctop + 0.55
    plot_right = cleft + cwidth - 0.1
    plot_bottom = ctop + cheight - 0.55  # reserve for category labels
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    max_val = max(values) if values else 0
    ticks, axis_top = _y_ticks(max_val)

    # Y axis labels + horizontal grid
    for tval in ticks:
        ty = plot_bottom - (tval / axis_top) * plot_h if axis_top > 0 else plot_bottom
        add_line(slide, plot_left, ty, plot_right, ty,
                 color=pal.grid_gray, width_pt=0.5)
        tb = add_textbox(slide, cleft, ty - 0.10, 0.5, 0.22,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"{int(round(tval))}",
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.RIGHT, first=True)

    # Bars
    n = len(values)
    slot_w = plot_w / n
    bar_w = slot_w * 0.6
    bar_centers = []
    for i, (cat, val) in enumerate(zip(categories, values)):
        slot_left = plot_left + i * slot_w
        bar_left = slot_left + (slot_w - bar_w) / 2
        bar_h = (val / axis_top) * plot_h if axis_top > 0 else 0
        bar_top = plot_bottom - bar_h
        # Color
        if forecast_from_index is not None and i >= forecast_from_index:
            fill = pal.bright_blue
        elif focus_index is not None and i == focus_index:
            fill = pal.bright_blue
        else:
            fill = pal.dark_navy
        add_rect(slide, bar_left, bar_top, bar_w, bar_h, fill=fill)

        # Value label above the bar
        tb = add_textbox(slide, bar_left - 0.2, bar_top - 0.30,
                         bar_w + 0.4, 0.25,
                         anchor=MSO_ANCHOR.BOTTOM)
        write_paragraph(tb.text_frame, f"{int(round(val))}",
                        size=typo.chart_label_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

        # Category label
        tb = add_textbox(slide, slot_left - 0.05, plot_bottom + 0.08,
                         slot_w + 0.1, 0.5,
                         anchor=MSO_ANCHOR.TOP)
        write_paragraph(tb.text_frame, str(cat),
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
        bar_centers.append(bar_left + bar_w / 2)

    return bar_centers, plot_bottom, bar_w, plot_top, axis_top


def _draw_growth_arrow(slide, theme, *, x1, y1, x2, y2, label_pct,
                       color=None):
    """Arrow with percentage in oval label centered on the arrow."""
    pal = theme.palette
    color = color or pal.text_dark

    from pptx.oxml.ns import qn
    from lxml import etree

    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                       Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.0)

    # Add arrow head end via XML
    ln = line.line._get_or_add_ln()
    for tail in ln.findall(qn("a:tailEnd")):
        ln.remove(tail)
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")

    # Oval label centered
    label_w, label_h = 0.55, 0.28
    midx = (x1 + x2) / 2 - label_w / 2
    midy = (y1 + y2) / 2 - label_h / 2
    add_oval(slide, midx, midy, label_w, label_h, fill=color)
    tb = add_textbox(slide, midx, midy, label_w, label_h,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, label_pct,
                    size=theme.typography.chart_label_size,
                    color=pal.white, family=theme.typography.family,
                    bold=False, align=PP_ALIGN.CENTER, first=True)


# ---------- Public slide builders ----------

def _common(prs, *, title, page_number, section_marker, source, footnote, theme,
            description_left, description_top, description_w,
            description="Description"):
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    _draw_description_header(slide, theme, left=description_left,
                              top=description_top, width=description_w,
                              label=description)
    return slide


def add_column_comparison(prs, *,
                          title="[Column chart for comparison with takeaways / Insert action title]",
                          categories: Sequence[str], values: Sequence[float],
                          focus_index: Optional[int] = None,
                          data_label="Data", data_unit="Unit",
                          takeaways: Sequence[str] = (),
                          description: str = "Description",
                          takeaway_header: str = "Key takeaways/main conclusion",
                          page_number=None, section_marker=None,
                          source="xx", footnote="1. xx",
                          theme: Theme = DEFAULT_THEME):
    slide = _common(prs, title=title, page_number=page_number,
                    section_marker=section_marker, source=source,
                    footnote=footnote, theme=theme,
                    description_left=DEFAULT_CHART_BOX[0],
                    description_top=DEFAULT_DESCRIPTION_TOP,
                    description_w=DEFAULT_CHART_BOX[2],
                    description=description)
    _draw_axis_and_bars(slide, theme, chart_box=DEFAULT_CHART_BOX,
                        data_label=data_label, data_unit=data_unit,
                        categories=categories, values=values,
                        focus_index=focus_index)
    _draw_takeaway(slide, theme, takeaways=takeaways, header=takeaway_header)
    return slide


def add_column_simple_growth(prs, *,
                             title="[Column chart with simple growth and takeaways / Insert action title]",
                             categories, values, growth_pct: str = "xx%",
                             data_label="Data", data_unit="Unit",
                             takeaways=(),
                             description: str = "Description",
                             takeaway_header: str = "Key takeaways/main conclusion",
                             page_number=None,
                             section_marker=None, source="xx",
                             footnote="1. xx",
                             theme: Theme = DEFAULT_THEME):
    slide = _common(prs, title=title, page_number=page_number,
                    section_marker=section_marker, source=source,
                    footnote=footnote, theme=theme,
                    description_left=DEFAULT_CHART_BOX[0],
                    description_top=DEFAULT_DESCRIPTION_TOP,
                    description_w=DEFAULT_CHART_BOX[2],
                    description=description)
    legend = [(theme.palette.dark_navy, "Actuals")]
    centers, baseline, bw, plot_top, axis_top = _draw_axis_and_bars(
        slide, theme, chart_box=DEFAULT_CHART_BOX,
        data_label=data_label, data_unit=data_unit,
        categories=categories, values=values, legend=legend,
    )
    # Single growth arrow ABOVE all data labels (which sit at bar_top - 0.30).
    if values:
        first_top = baseline - (values[0] / axis_top) * (baseline - plot_top)
        last_top = baseline - (values[-1] / axis_top) * (baseline - plot_top)
        _draw_growth_arrow(slide, theme,
                           x1=centers[0], y1=first_top - 0.45,
                           x2=centers[-1] + 0.15, y2=last_top - 0.55,
                           label_pct=growth_pct, color=theme.palette.dark_navy)
    _draw_takeaway(slide, theme, takeaways=takeaways, header=takeaway_header)
    return slide


def add_column_split_growth(prs, *,
                            title="[Column chart with split growth and takeaways / Insert action title]",
                            categories, values, split_index: int,
                            growth_pct_first: str = "xx%",
                            growth_pct_second: str = "xx%",
                            data_label="Data", data_unit="Unit",
                            takeaways=(),
                            description: str = "Description",
                            takeaway_header: str = "Key takeaways/main conclusion",
                            page_number=None,
                            section_marker=None, source="xx",
                            footnote="1. xx",
                            theme: Theme = DEFAULT_THEME):
    slide = _common(prs, title=title, page_number=page_number,
                    section_marker=section_marker, source=source,
                    footnote=footnote, theme=theme,
                    description_left=DEFAULT_CHART_BOX[0],
                    description_top=DEFAULT_DESCRIPTION_TOP,
                    description_w=DEFAULT_CHART_BOX[2],
                    description=description)
    legend = [(theme.palette.dark_navy, "Actuals")]
    centers, baseline, bw, plot_top, axis_top = _draw_axis_and_bars(
        slide, theme, chart_box=DEFAULT_CHART_BOX,
        data_label=data_label, data_unit=data_unit,
        categories=categories, values=values, legend=legend,
    )
    if values and split_index > 0:
        def _y(v):
            return baseline - (v / axis_top) * (baseline - plot_top)
        # First arrow (sits above the data labels in its segment)
        _draw_growth_arrow(slide, theme,
                           x1=centers[0], y1=_y(values[0]) - 0.45,
                           x2=centers[split_index] + 0.10,
                           y2=_y(values[split_index]) - 0.55,
                           label_pct=growth_pct_first,
                           color=theme.palette.dark_navy)
        # Second arrow
        _draw_growth_arrow(slide, theme,
                           x1=centers[split_index],
                           y1=_y(values[split_index]) - 0.55,
                           x2=centers[-1] + 0.10,
                           y2=_y(values[-1]) - 0.55,
                           label_pct=growth_pct_second,
                           color=theme.palette.dark_navy)
    _draw_takeaway(slide, theme, takeaways=takeaways, header=takeaway_header)
    return slide


def add_column_historic_forecast(prs, *,
                                 title="[Column chart with historic and forecast figures / Insert action title]",
                                 categories, values, forecast_from_index: int,
                                 historic_growth: str = "xx%",
                                 forecast_growth: str = "xx%",
                                 data_label="Data", data_unit="Unit",
                                 takeaways=(),
                                 description: str = "Description",
                                 takeaway_header: str = "Key takeaways/main conclusion",
                                 page_number=None,
                                 section_marker=None, source="xx",
                                 footnote="1. xx",
                                 theme: Theme = DEFAULT_THEME):
    slide = _common(prs, title=title, page_number=page_number,
                    section_marker=section_marker, source=source,
                    footnote=footnote, theme=theme,
                    description_left=DEFAULT_CHART_BOX[0],
                    description_top=DEFAULT_DESCRIPTION_TOP,
                    description_w=DEFAULT_CHART_BOX[2],
                    description=description)
    legend = [(theme.palette.dark_navy, "Actuals"),
              (theme.palette.bright_blue, "Forecast")]
    centers, baseline, bw, plot_top, axis_top = _draw_axis_and_bars(
        slide, theme, chart_box=DEFAULT_CHART_BOX,
        data_label=data_label, data_unit=data_unit,
        categories=categories, values=values, legend=legend,
        forecast_from_index=forecast_from_index,
    )
    if values and 0 < forecast_from_index < len(values):
        def _y(v):
            return baseline - (v / axis_top) * (baseline - plot_top)
        # Historic arrow (dark) — clears value labels of historic bars
        _draw_growth_arrow(slide, theme,
                           x1=centers[0], y1=_y(values[0]) - 0.45,
                           x2=centers[forecast_from_index - 1] + 0.10,
                           y2=_y(values[forecast_from_index - 1]) - 0.55,
                           label_pct=historic_growth,
                           color=theme.palette.dark_navy)
        # Forecast arrow (light blue)
        _draw_growth_arrow(slide, theme,
                           x1=centers[forecast_from_index - 1],
                           y1=_y(values[forecast_from_index - 1]) - 0.55,
                           x2=centers[-1] + 0.10,
                           y2=_y(values[-1]) - 0.55,
                           label_pct=forecast_growth,
                           color=theme.palette.bright_blue)
    _draw_takeaway(slide, theme, takeaways=takeaways, header=takeaway_header)
    return slide
