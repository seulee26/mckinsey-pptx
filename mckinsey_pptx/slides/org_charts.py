"""Org / team / hierarchy slides:
 - issue_tree: root issue -> main drivers -> secondary -> underlying
 - org_chart: multi-level boxed hierarchy with connector lines
 - project_team_circles: leader circle on top + N teammate circles below
 - team_chart: project-team header + N function columns + role tiles per column
"""
from __future__ import annotations
from typing import Sequence, Optional, Dict, List, Tuple

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph, add_subtitle_placeholder,
)
from ..theme import Theme, DEFAULT_THEME


# ---------- shared primitives ----------

def _box(slide, theme, left, top, w, h, *, text, fill, color, bold=False,
         size=None):
    typo = theme.typography
    add_rect(slide, left, top, w, h, fill=fill)
    tb = add_textbox(slide, left + 0.05, top, w - 0.10, h,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, text,
                    size=size or typo.body_size, bold=bold,
                    color=color, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)


def _elbow(slide, theme, x1, y1, x2, y2):
    """Draw a 2-segment elbow connector from (x1,y1) horizontal-then-vertical to (x2,y2)."""
    pal = theme.palette
    midx = (x1 + x2) / 2
    add_line(slide, x1, y1, midx, y1, color=pal.rule_gray, width_pt=0.6)
    add_line(slide, midx, y1, midx, y2, color=pal.rule_gray, width_pt=0.6)
    add_line(slide, midx, y2, x2, y2, color=pal.rule_gray, width_pt=0.6)


# ---------- Issue tree ----------

def add_issue_tree(prs, *,
                   title="[Issue tree / Insert action title]",
                   subtitle: Optional[str] = "[Insert subtitle]",
                   root: str = "[Main issue]",
                   main_drivers: Sequence[Dict],
                   page_number=None, section_marker="Section marker",
                   source="xx", footnote="1. xx",
                   theme: Theme = DEFAULT_THEME):
    """main_drivers: [{
        "label": "[Main drivers of issue]",
        "secondaries": [
            {"label": "[Secondary drivers of issue]",
             "underlying": ["[Underlying drivers]", "[Underlying drivers]"]}
        ]
    }]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    body_top = 1.95
    body_bottom = layout.footer_top_in - 0.15
    body_h = body_bottom - body_top

    # Column geometry
    col_w_root = 1.4
    col_w_main = 1.6
    col_w_sec = 2.0
    col_w_und = 2.2
    gap = 0.35
    left = layout.margin_left_in
    x_root = left
    x_main = x_root + col_w_root + gap
    x_sec = x_main + col_w_main + gap
    x_und = x_sec + col_w_sec + gap

    # Root box (big, deep navy)
    root_h = 1.6
    root_y = body_top + (body_h - root_h) / 2
    _box(slide, theme, x_root, root_y, col_w_root, root_h,
         text=root, fill=pal.deep_navy, color=pal.white, bold=True,
         size=typo.body_size)

    # Underlying drivers count drives layout
    total_underlying = sum(
        max(1, len(s.get("underlying", []) or ["[Underlying drivers]"]))
        for m in main_drivers
        for s in m.get("secondaries", []) or [{"underlying": []}]
    )
    if total_underlying == 0:
        total_underlying = 1

    und_h = min(0.55, body_h / max(total_underlying, 1) - 0.08)
    und_h = max(und_h, 0.30)
    und_gap = max(0.06, (body_h - und_h * total_underlying) / max(total_underlying, 1))

    y_cursor = body_top
    main_root_y = root_y + root_h / 2

    for m in main_drivers:
        secs = m.get("secondaries", []) or []
        # compute total underlying for this main
        m_total_und = sum(max(1, len(s.get("underlying", []) or [])) for s in secs)
        if m_total_und == 0:
            m_total_und = 1
        m_height = m_total_und * (und_h + und_gap)
        m_box_h = min(0.7, m_height * 0.55)
        m_box_h = max(m_box_h, 0.45)
        m_box_y = y_cursor + (m_height - m_box_h) / 2

        _box(slide, theme, x_main, m_box_y, col_w_main, m_box_h,
             text=m.get("label", "[Main drivers of issue]"),
             fill=pal.bright_blue, color=pal.white, bold=True,
             size=typo.body_size - 1)
        # connector root -> main
        _elbow(slide, theme,
               x_root + col_w_root, main_root_y,
               x_main, m_box_y + m_box_h / 2)

        sec_y = y_cursor
        for s in secs:
            unds = s.get("underlying", []) or ["[Underlying drivers]"]
            s_height = max(1, len(unds)) * (und_h + und_gap)
            s_box_h = min(0.6, s_height * 0.6)
            s_box_h = max(s_box_h, 0.40)
            s_box_y = sec_y + (s_height - s_box_h) / 2

            _box(slide, theme, x_sec, s_box_y, col_w_sec, s_box_h,
                 text=s.get("label", "[Secondary drivers of issue]"),
                 fill=pal.soft_gray, color=pal.text_dark, bold=False,
                 size=typo.body_size - 1)
            _elbow(slide, theme,
                   x_main + col_w_main, m_box_y + m_box_h / 2,
                   x_sec, s_box_y + s_box_h / 2)

            uy = sec_y
            for u in unds:
                _box(slide, theme, x_und, uy, col_w_und, und_h,
                     text=u, fill=pal.soft_gray, color=pal.text_dark,
                     bold=False, size=typo.body_size - 2)
                _elbow(slide, theme,
                       x_sec + col_w_sec, s_box_y + s_box_h / 2,
                       x_und, uy + und_h / 2)
                uy += und_h + und_gap
            sec_y += s_height
        y_cursor += m_height
    return slide


# ---------- Organizational chart ----------

def add_org_chart(prs, *,
                  title="[Organizational chart / Insert action title]",
                  subtitle: Optional[str] = "[Insert subtitle]",
                  ceo: str = "[Name/title/role]",
                  branches: Sequence[Dict],
                  page_number=None, section_marker="Section marker",
                  source="xx", footnote="1. xx",
                  theme: Theme = DEFAULT_THEME):
    """branches: [{"head": "[Name/title/role]", "reports": ["[Name/title/role]", ...]}]
    Maximum 5-6 branches recommended.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    body_top = 1.95
    body_bottom = layout.footer_top_in - 0.15
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    box_w = 1.85
    box_h = 0.45

    # CEO at top center
    ceo_x = layout.margin_left_in + (width - box_w) / 2
    ceo_y = body_top + 0.10
    _box(slide, theme, ceo_x, ceo_y, box_w, box_h, text=ceo,
         fill=pal.deep_navy, color=pal.white, bold=True,
         size=typo.body_size - 1)

    # Heads row
    n = len(branches)
    if n == 0:
        return slide
    head_y = ceo_y + box_h + 0.6
    total_head_w = n * box_w + (n - 1) * 0.25
    head_left0 = layout.margin_left_in + (width - total_head_w) / 2
    head_xs = [head_left0 + i * (box_w + 0.25) for i in range(n)]

    # Vertical line from CEO down + horizontal connector
    spine_y = ceo_y + box_h + 0.30
    add_line(slide, ceo_x + box_w / 2, ceo_y + box_h,
             ceo_x + box_w / 2, spine_y, color=pal.rule_gray, width_pt=0.6)
    if n > 1:
        add_line(slide, head_xs[0] + box_w / 2, spine_y,
                 head_xs[-1] + box_w / 2, spine_y,
                 color=pal.rule_gray, width_pt=0.6)
    for hx in head_xs:
        add_line(slide, hx + box_w / 2, spine_y,
                 hx + box_w / 2, head_y, color=pal.rule_gray, width_pt=0.6)

    # Heads + reports
    reports_top = head_y + box_h + 0.30
    avail_h = body_bottom - reports_top
    for hx, br in zip(head_xs, branches):
        _box(slide, theme, hx, head_y, box_w, box_h,
             text=br.get("head", "[Name/title/role]"),
             fill=pal.mid_blue, color=pal.white, bold=False,
             size=typo.body_size - 2)
        reports = br.get("reports", []) or []
        if not reports:
            continue
        rh = min(0.42, avail_h / len(reports) - 0.05)
        rh = max(rh, 0.30)
        rg = 0.10
        ry = reports_top
        # spine from head down
        add_line(slide, hx + box_w / 2, head_y + box_h,
                 hx + box_w / 2, reports_top + rh / 2 + (len(reports) - 1)
                 * (rh + rg) - 0.0,
                 color=pal.rule_gray, width_pt=0.5)
        for r in reports:
            _box(slide, theme, hx, ry, box_w, rh, text=r,
                 fill=pal.soft_gray, color=pal.text_dark,
                 bold=False, size=typo.body_size - 3)
            # Tee from spine
            add_line(slide, hx + box_w / 2, ry + rh / 2,
                     hx, ry + rh / 2,
                     color=pal.rule_gray, width_pt=0.5)
            ry += rh + rg
    return slide


# ---------- Project team circles ----------

def _circle_with_icon(slide, theme, cx, cy, d, *, fill, icon: str,
                       icon_color=None, font_size=20):
    pal, typo = theme.palette, theme.typography
    add_oval(slide, cx - d / 2, cy - d / 2, d, d, fill=fill)
    tb = add_textbox(slide, cx - d / 2, cy - d / 2, d, d,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, icon, size=font_size,
                    color=icon_color or pal.white, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)


def add_project_team_circles(prs, *,
                             title="[Project team or functions / Insert action title]",
                             subtitle: Optional[str] = "[Insert subtitle]",
                             leader: Dict,
                             members: Sequence[Dict],
                             page_number=None, section_marker="Section marker",
                             source="xx", footnote="1. xx",
                             theme: Theme = DEFAULT_THEME):
    """leader: {"name": "[Insert name/title]", "description": "...", "icon": "👥"}
    members: same shape; up to ~5 nodes.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 2.0

    # Leader circle (centered)
    leader_d = 1.0
    leader_cx = layout.margin_left_in + width / 2
    leader_cy = body_top + leader_d / 2
    _circle_with_icon(slide, theme, leader_cx, leader_cy, leader_d,
                       fill=pal.deep_navy, icon=leader.get("icon", "★"),
                       font_size=22)
    # Leader label/description to the right
    tb = add_textbox(slide, leader_cx + leader_d / 2 + 0.15,
                     leader_cy - 0.30, 3.5, 0.30)
    write_paragraph(tb.text_frame, leader.get("name", "[Insert name/title]"),
                    size=typo.body_size, bold=True, color=pal.text_dark,
                    family=typo.family, first=True)
    if leader.get("description"):
        tb = add_textbox(slide, leader_cx + leader_d / 2 + 0.15,
                         leader_cy - 0.02, 3.8, 0.6)
        write_paragraph(tb.text_frame, leader["description"],
                        size=typo.body_size - 1, color=pal.text_dark,
                        family=typo.family, first=True)

    # Connector line from leader down
    spine_y = leader_cy + leader_d / 2 + 0.50
    add_line(slide, leader_cx, leader_cy + leader_d / 2,
             leader_cx, spine_y, color=pal.rule_gray, width_pt=0.6)

    # Member circles in a row
    n = len(members)
    if n == 0:
        return slide
    member_d = 0.80
    member_cy = spine_y + 0.35 + member_d / 2
    avail = width - 0.4
    slot = avail / n
    member_cx_start = layout.margin_left_in + 0.2 + slot / 2

    if n > 1:
        add_line(slide, member_cx_start, spine_y,
                 member_cx_start + (n - 1) * slot, spine_y,
                 color=pal.rule_gray, width_pt=0.6)
    for i, m in enumerate(members):
        cx = member_cx_start + i * slot
        # Drop from spine
        add_line(slide, cx, spine_y, cx, member_cy - member_d / 2,
                 color=pal.rule_gray, width_pt=0.6)
        _circle_with_icon(slide, theme, cx, member_cy, member_d,
                           fill=pal.soft_gray, icon=m.get("icon", "•"),
                           icon_color=pal.deep_navy, font_size=18)
        tb = add_textbox(slide, cx - slot / 2 + 0.05,
                         member_cy + member_d / 2 + 0.10,
                         slot - 0.10, 0.32)
        write_paragraph(tb.text_frame, m.get("name", "[Insert name/title]"),
                        size=typo.body_size, bold=True, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER,
                        first=True)
        if m.get("description"):
            tb = add_textbox(slide, cx - slot / 2 + 0.05,
                             member_cy + member_d / 2 + 0.42,
                             slot - 0.10, 1.1)
            write_paragraph(tb.text_frame, m["description"],
                            size=typo.body_size - 2,
                            color=pal.text_dark, family=typo.family,
                            align=PP_ALIGN.CENTER, first=True)
    return slide


# ---------- Team chart (function columns) ----------

def add_team_chart(prs, *,
                   title="[Team chart / Insert action title]",
                   project_name: str = "[Project] Team",
                   functions: Sequence[Dict],
                   page_number=None, section_marker="Section marker",
                   source="xx", footnote="1. xx",
                   theme: Theme = DEFAULT_THEME):
    """functions: [{"name": "[Insert function]", "description": "...",
                    "roles": [{"name": "[Insert role]", "kind": "filled"|"outline"}, ...]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.55

    # Top circle (project team)
    bub_d = 0.85
    bub_cx = layout.margin_left_in + width / 2 - 2.0
    bub_cy = body_top + bub_d / 2
    add_oval(slide, bub_cx - bub_d / 2, bub_cy - bub_d / 2, bub_d, bub_d,
             fill=pal.bright_blue)
    tb = add_textbox(slide, bub_cx - bub_d / 2, bub_cy - bub_d / 2,
                     bub_d, bub_d, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, project_name, size=typo.body_size - 1,
                    bold=True, color=pal.white, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    # Legend (top-right) — filled vs outline circles
    leg_y = body_top + 0.10
    leg_x = layout.margin_left_in + width - 4.6
    add_oval(slide, leg_x, leg_y, 0.30, 0.30, fill=None,
             line=pal.deep_navy, line_width=0.8)
    tb = add_textbox(slide, leg_x + 0.40, leg_y - 0.02, 2.0, 0.32,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "[Type of role, if relevant]",
                    size=typo.chart_label_size, color=pal.text_dark,
                    family=typo.family, first=True)
    leg_x2 = leg_x + 2.40
    add_oval(slide, leg_x2, leg_y, 0.30, 0.30, fill=pal.deep_navy)
    tb = add_textbox(slide, leg_x2 + 0.40, leg_y - 0.02, 2.0, 0.32,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "[Type of role, if relevant]",
                    size=typo.chart_label_size, color=pal.text_dark,
                    family=typo.family, first=True)

    # Function columns
    cols_top = body_top + bub_d + 0.45
    cols_bottom = layout.footer_top_in - 0.20
    n = len(functions)
    if n == 0:
        return slide
    col_w = (width - (n - 1) * 0.05) / n
    for i, fn in enumerate(functions):
        cx = layout.margin_left_in + i * (col_w + 0.05)
        # Header band (gray)
        head_h = 1.0
        add_rect(slide, cx, cols_top, col_w, head_h, fill=pal.soft_gray)
        tb = add_textbox(slide, cx + 0.10, cols_top + 0.05, col_w - 0.20, 0.32)
        write_paragraph(tb.text_frame, fn.get("name", "[Insert function]"),
                        size=typo.body_size - 1, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        if fn.get("description"):
            tb = add_textbox(slide, cx + 0.10, cols_top + 0.40,
                             col_w - 0.20, head_h - 0.45)
            write_paragraph(tb.text_frame, fn["description"],
                            size=typo.body_size - 2, color=pal.text_dark,
                            family=typo.family, first=True)
        # Connector from project bubble down to first column header center?
        # The reference shows light vertical lines from top to each col head — skip for compactness.

        # Roles
        ry = cols_top + head_h + 0.15
        rd = 0.30
        for role in fn.get("roles", []) or []:
            kind = role.get("kind", "filled")
            if kind == "filled":
                add_oval(slide, cx + 0.10, ry, rd, rd, fill=pal.deep_navy)
                color = pal.white
            else:
                add_oval(slide, cx + 0.10, ry, rd, rd, fill=None,
                         line=pal.deep_navy, line_width=0.8)
                color = pal.deep_navy
            tb = add_textbox(slide, cx + 0.10, ry, rd, rd,
                             anchor=MSO_ANCHOR.MIDDLE)
            # Person icon glyph
            write_paragraph(tb.text_frame, "♟",
                            size=12, color=color,
                            family=typo.family, align=PP_ALIGN.CENTER,
                            first=True)
            tb = add_textbox(slide, cx + 0.10 + rd + 0.08, ry,
                             col_w - rd - 0.30, rd,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, role.get("name", "[Insert role]"),
                            size=typo.body_size - 2, color=pal.text_dark,
                            family=typo.family, first=True)
            ry += rd + 0.06
    return slide
