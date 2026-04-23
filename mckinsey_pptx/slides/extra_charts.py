"""Additional chart types: stacked column, grouped column, line chart.
All share the McKinsey design system (deep navy / bright blue palette,
[Description] header, takeaway pane on the right).
"""
from __future__ import annotations
from typing import Sequence, Optional, List, Dict, Tuple
import math

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME
from .column_chart import (
    DEFAULT_CHART_BOX, DEFAULT_TAKEAWAY_BOX, DEFAULT_DESCRIPTION_TOP,
    _draw_takeaway, _draw_description_header, _y_ticks,
)


def _palette_series(theme, n):
    pal = theme.palette
    base = [pal.deep_navy, pal.bright_blue, pal.mid_blue, pal.light_blue,
            pal.royal_blue, pal.placeholder_gray, pal.status_amber,
            pal.status_green]
    return [base[i % len(base)] for i in range(n)]


def _draw_axis_frame(slide, theme, *, chart_box, data_label, data_unit,
                     ticks, axis_top, legend=None):
    pal, typo = theme.palette, theme.typography
    cleft, ctop, cwidth, cheight = chart_box

    # [Data], [Unit] header
    tb = add_textbox(slide, cleft, ctop, cwidth - 1.5, 0.30)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = f"[{data_label}], "
    r.font.size = Pt(typo.section_title_size); r.font.bold = True
    r.font.color.rgb = pal.text_dark; r.font.name = typo.family
    r2 = p.add_run(); r2.text = f"[{data_unit}]"
    r2.font.size = Pt(typo.section_title_size)
    r2.font.color.rgb = pal.placeholder_gray; r2.font.name = typo.family

    # Legend (top-right)
    if legend:
        leg_y = ctop
        leg_left = cleft + cwidth - 0.05
        for color_rgb, label in reversed(legend):
            label_w = max(0.7, len(label) * 0.08)
            swatch_w = 0.22
            tb = add_textbox(slide, leg_left - label_w, leg_y, label_w, 0.25,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, label, size=typo.chart_label_size,
                            color=pal.text_dark, family=typo.family,
                            align=PP_ALIGN.LEFT, first=True)
            add_rect(slide, leg_left - label_w - swatch_w - 0.05,
                     leg_y + 0.05, swatch_w, 0.15, fill=color_rgb)
            leg_left -= label_w + swatch_w + 0.20

    # Plot area
    plot_left = cleft + 0.55
    plot_top = ctop + 0.55
    plot_right = cleft + cwidth - 0.1
    plot_bottom = ctop + cheight - 0.55
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    # Y-axis labels + grid
    for tval in ticks:
        ty = plot_bottom - (tval / axis_top) * plot_h if axis_top > 0 else plot_bottom
        add_line(slide, plot_left, ty, plot_right, ty,
                 color=pal.grid_gray, width_pt=0.5)
        tb = add_textbox(slide, cleft, ty - 0.10, 0.5, 0.22,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"{int(round(tval))}",
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.RIGHT, first=True)
    return plot_left, plot_top, plot_right, plot_bottom, plot_w, plot_h


# ---------- Stacked column chart ----------

def add_stacked_column_chart(prs, *,
                              title="[Stacked column / Insert action title]",
                              categories: Sequence[str],
                              series: Sequence[Dict],
                              data_label="Data", data_unit="Unit",
                              takeaways: Sequence[str] = (),
                              show_totals: bool = True,
                              page_number=None, section_marker=None,
                              source="xx", footnote="1. xx",
                              theme: Theme = DEFAULT_THEME):
    """series: [{"name": "Segment A", "values": [v1, v2, v3, ...]}]
    Each series has one value per category. Bars are stacked.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    _draw_description_header(slide, theme, left=DEFAULT_CHART_BOX[0],
                              top=DEFAULT_DESCRIPTION_TOP,
                              width=DEFAULT_CHART_BOX[2])
    chart_box = DEFAULT_CHART_BOX

    # totals per category
    totals = [sum(s["values"][i] if i < len(s["values"]) else 0
                   for s in series) for i in range(len(categories))]
    max_total = max(totals) if totals else 0
    ticks, axis_top = _y_ticks(max_total)

    colors = _palette_series(theme, len(series))
    legend = list(zip(colors, [s.get("name", f"Series {i+1}")
                                for i, s in enumerate(series)]))

    pl, pt, pr, pb, pw, ph = _draw_axis_frame(
        slide, theme, chart_box=chart_box,
        data_label=data_label, data_unit=data_unit,
        ticks=ticks, axis_top=axis_top, legend=legend)

    n = len(categories)
    if n == 0 or axis_top <= 0:
        _draw_takeaway(slide, theme, takeaways=takeaways)
        return slide
    slot_w = pw / n
    bar_w = slot_w * 0.55

    for i, cat in enumerate(categories):
        slot_left = pl + i * slot_w
        bar_left = slot_left + (slot_w - bar_w) / 2
        # Stack bottom-up
        cum = 0
        for s_idx, s in enumerate(series):
            v = s["values"][i] if i < len(s["values"]) else 0
            if v <= 0:
                continue
            seg_h = (v / axis_top) * ph
            seg_top = pb - (cum + v) / axis_top * ph
            add_rect(slide, bar_left, seg_top, bar_w, seg_h,
                     fill=colors[s_idx])
            # In-segment label if segment is tall enough
            if seg_h > 0.30:
                tb = add_textbox(slide, bar_left, seg_top, bar_w, seg_h,
                                 anchor=MSO_ANCHOR.MIDDLE)
                write_paragraph(tb.text_frame, f"{int(round(v))}",
                                size=typo.chart_label_size, color=pal.white,
                                family=typo.family,
                                align=PP_ALIGN.CENTER, first=True)
            cum += v
        # Total above bar
        if show_totals:
            top_y = pb - (cum / axis_top) * ph
            tb = add_textbox(slide, bar_left - 0.2, top_y - 0.30,
                             bar_w + 0.4, 0.25, anchor=MSO_ANCHOR.BOTTOM)
            write_paragraph(tb.text_frame, f"{int(round(cum))}",
                            size=typo.chart_label_size, bold=True,
                            color=pal.text_dark, family=typo.family,
                            align=PP_ALIGN.CENTER, first=True)
        # Category label
        tb = add_textbox(slide, slot_left - 0.05, pb + 0.08,
                         slot_w + 0.1, 0.5, anchor=MSO_ANCHOR.TOP)
        write_paragraph(tb.text_frame, str(cat),
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    _draw_takeaway(slide, theme, takeaways=takeaways)
    return slide


# ---------- Grouped column chart ----------

def add_grouped_column_chart(prs, *,
                              title="[Grouped column / Insert action title]",
                              categories: Sequence[str],
                              series: Sequence[Dict],
                              data_label="Data", data_unit="Unit",
                              takeaways: Sequence[str] = (),
                              show_values: bool = True,
                              page_number=None, section_marker=None,
                              source="xx", footnote="1. xx",
                              theme: Theme = DEFAULT_THEME):
    """series: [{"name": "2023", "values": [v1, v2, ...]}, ...]
    Each series renders as one bar within each category group.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    _draw_description_header(slide, theme, left=DEFAULT_CHART_BOX[0],
                              top=DEFAULT_DESCRIPTION_TOP,
                              width=DEFAULT_CHART_BOX[2])
    chart_box = DEFAULT_CHART_BOX

    max_val = 0
    for s in series:
        for v in s.get("values", []):
            if v > max_val: max_val = v
    ticks, axis_top = _y_ticks(max_val)

    colors = _palette_series(theme, len(series))
    legend = list(zip(colors, [s.get("name", f"Series {i+1}")
                                for i, s in enumerate(series)]))

    pl, pt, pr, pb, pw, ph = _draw_axis_frame(
        slide, theme, chart_box=chart_box,
        data_label=data_label, data_unit=data_unit,
        ticks=ticks, axis_top=axis_top, legend=legend)

    n_cat = len(categories)
    n_ser = max(len(series), 1)
    if n_cat == 0 or axis_top <= 0:
        _draw_takeaway(slide, theme, takeaways=takeaways)
        return slide
    slot_w = pw / n_cat
    group_w = slot_w * 0.75
    bar_w = group_w / n_ser
    inner_pad = bar_w * 0.10

    for ci, cat in enumerate(categories):
        slot_left = pl + ci * slot_w
        group_left = slot_left + (slot_w - group_w) / 2
        for si, s in enumerate(series):
            v = s["values"][ci] if ci < len(s.get("values", [])) else 0
            bar_left = group_left + si * bar_w + inner_pad / 2
            bw = bar_w - inner_pad
            bh = (v / axis_top) * ph
            bar_top = pb - bh
            add_rect(slide, bar_left, bar_top, bw, bh, fill=colors[si])
            if show_values and bh > 0.05:
                tb = add_textbox(slide, bar_left - 0.15, bar_top - 0.28,
                                 bw + 0.30, 0.25, anchor=MSO_ANCHOR.BOTTOM)
                write_paragraph(tb.text_frame, f"{int(round(v))}",
                                size=typo.chart_label_size,
                                color=pal.text_dark, family=typo.family,
                                align=PP_ALIGN.CENTER, first=True)
        # Category label
        tb = add_textbox(slide, slot_left - 0.05, pb + 0.08,
                         slot_w + 0.1, 0.5, anchor=MSO_ANCHOR.TOP)
        write_paragraph(tb.text_frame, str(cat),
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    _draw_takeaway(slide, theme, takeaways=takeaways)
    return slide


# ---------- Line chart ----------

def add_line_chart(prs, *,
                   title="[Line chart / Insert action title]",
                   categories: Sequence,
                   series: Sequence[Dict],
                   data_label="Data", data_unit="Unit",
                   takeaways: Sequence[str] = (),
                   show_markers: bool = True,
                   show_values_for: Optional[Sequence[str]] = None,
                   page_number=None, section_marker=None,
                   source="xx", footnote="1. xx",
                   theme: Theme = DEFAULT_THEME):
    """series: [{"name": "US", "values": [v1, v2, ...], "color": optional}]
    Up to ~4 lines recommended.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    _draw_description_header(slide, theme, left=DEFAULT_CHART_BOX[0],
                              top=DEFAULT_DESCRIPTION_TOP,
                              width=DEFAULT_CHART_BOX[2])
    chart_box = DEFAULT_CHART_BOX

    color_map = {
        "navy": pal.deep_navy, "blue": pal.bright_blue,
        "mid_blue": pal.mid_blue, "light_blue": pal.light_blue,
        "royal": pal.royal_blue, "amber": pal.status_amber,
        "green": pal.status_green, "red": pal.status_red,
    }
    default_palette = [pal.deep_navy, pal.bright_blue, pal.mid_blue,
                       pal.royal_blue, pal.status_amber]
    series_colors = []
    for i, s in enumerate(series):
        c = s.get("color")
        if isinstance(c, str) and c in color_map:
            series_colors.append(color_map[c])
        else:
            series_colors.append(default_palette[i % len(default_palette)])

    # Determine y range
    all_vals = [v for s in series for v in s.get("values", [])]
    max_val = max(all_vals) if all_vals else 0
    min_val = min(all_vals) if all_vals else 0
    if min_val < 0:
        ticks, axis_top = _y_ticks(max(abs(max_val), abs(min_val)))
        # symmetric? simpler: keep positive only for now
    else:
        ticks, axis_top = _y_ticks(max_val)

    legend = list(zip(series_colors,
                      [s.get("name", f"Series {i+1}")
                       for i, s in enumerate(series)]))

    pl, pt, pr, pb, pw, ph = _draw_axis_frame(
        slide, theme, chart_box=chart_box,
        data_label=data_label, data_unit=data_unit,
        ticks=ticks, axis_top=axis_top, legend=legend)

    n_pts = len(categories)
    if n_pts == 0 or axis_top <= 0:
        _draw_takeaway(slide, theme, takeaways=takeaways)
        return slide
    if n_pts == 1:
        x_step = 0
        x_origin = pl + pw / 2
    else:
        x_step = pw / (n_pts - 1)
        x_origin = pl

    show_value_set = set(show_values_for or [])

    # Draw lines + markers
    for si, s in enumerate(series):
        color = series_colors[si]
        vals = s.get("values", [])
        pts = []
        for i, v in enumerate(vals[:n_pts]):
            x = x_origin + i * x_step
            y = pb - (v / axis_top) * ph if axis_top > 0 else pb
            pts.append((x, y))
        # connect with line segments
        for (x1, y1), (x2, y2) in zip(pts[:-1], pts[1:]):
            add_line(slide, x1, y1, x2, y2, color=color, width_pt=2.0)
        # markers + optional value labels
        if show_markers:
            for x, y in pts:
                d = 0.14
                add_oval(slide, x - d / 2, y - d / 2, d, d, fill=color,
                         line=pal.white, line_width=0.8)
        if s.get("name") in show_value_set:
            for (x, y), v in zip(pts, vals[:n_pts]):
                tb = add_textbox(slide, x - 0.30, y - 0.36, 0.60, 0.22,
                                 anchor=MSO_ANCHOR.MIDDLE)
                write_paragraph(tb.text_frame, f"{v:g}",
                                size=typo.chart_label_size, color=color,
                                family=typo.family, align=PP_ALIGN.CENTER,
                                first=True)

    # X-axis labels
    for i, cat in enumerate(categories):
        x = x_origin + i * x_step
        tb = add_textbox(slide, x - 0.5, pb + 0.08, 1.0, 0.30,
                         anchor=MSO_ANCHOR.TOP)
        write_paragraph(tb.text_frame, str(cat),
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    _draw_takeaway(slide, theme, takeaways=takeaways)
    return slide
