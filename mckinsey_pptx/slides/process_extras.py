"""Process / metric overview slides:
 - process_flow_horizontal: 4-6 step flow with chevron-style arrows
 - funnel: top-down funnel (TAM/SAM/SOM, marketing funnel, etc.)
 - kpi_dashboard: 4-6 KPI tiles with big number + label + delta
"""
from __future__ import annotations
from typing import Sequence, Optional, Dict, Literal

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph, add_subtitle_placeholder,
)
from ..theme import Theme, DEFAULT_THEME


# ---------- Process flow (horizontal) ----------

def add_process_flow_horizontal(prs, *,
                                title="[Process flow / Insert action title]",
                                subtitle: Optional[str] = None,
                                steps: Sequence[Dict],
                                page_number=None,
                                section_marker="Section marker",
                                source="xx", footnote="1. xx",
                                theme: Theme = DEFAULT_THEME):
    """steps: [{"name": "Step 1", "description": "..."}]
    Renders 4-6 numbered chevron tiles in a row with descriptions below.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, 1.30,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.32)
        write_paragraph(tb.text_frame, subtitle, size=typo.body_size,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 2.0
    body_bottom = layout.footer_top_in - 0.30
    n = max(len(steps), 1)

    chev_h = 0.85
    overlap = 0.30
    chev_w = (width + (n - 1) * overlap) / n
    chev_y = body_top
    # Alternate fills for visual rhythm
    fills = [pal.deep_navy, pal.bright_blue, pal.deep_navy, pal.bright_blue,
             pal.deep_navy, pal.bright_blue]

    for i, st in enumerate(steps):
        cx = layout.margin_left_in + i * (chev_w - overlap)
        fill = fills[i % len(fills)]
        s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                   Inches(cx), Inches(chev_y),
                                   Inches(chev_w), Inches(chev_h))
        s.shadow.inherit = False
        s.fill.solid(); s.fill.fore_color.rgb = fill
        s.line.fill.background()

        # Number circle inside chevron
        n_d = 0.42
        nx = cx + 0.18
        ny = chev_y + (chev_h - n_d) / 2
        add_oval(slide, nx, ny, n_d, n_d, fill=pal.white)
        tb = add_textbox(slide, nx, ny, n_d, n_d, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(i + 1).zfill(2),
                        size=typo.body_size - 1, bold=True,
                        color=fill, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

        # Step name (overlay text)
        tb = add_textbox(slide, nx + n_d + 0.10, chev_y,
                         chev_w - n_d - 0.55, chev_h, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, st.get("name", f"Step {i+1}"),
                        size=typo.body_size, bold=True, color=pal.white,
                        family=typo.family, first=True)

        # Description below
        desc_top = chev_y + chev_h + 0.20
        tb = add_textbox(slide, cx + 0.10, desc_top,
                         chev_w - overlap - 0.10, body_bottom - desc_top - 0.10)
        write_paragraph(tb.text_frame, st.get("description", ""),
                        size=typo.body_size - 1, color=pal.text_dark,
                        family=typo.family, first=True)
    return slide


# ---------- Funnel ----------

def add_funnel(prs, *,
               title="[Funnel / Insert action title]",
               subtitle: Optional[str] = None,
               stages: Sequence[Dict],
               page_number=None,
               section_marker="Section marker",
               source="xx", footnote="1. xx",
               theme: Theme = DEFAULT_THEME):
    """stages: [{"name": "Awareness", "value": "1.2M", "description": "..."}]
    Top-down funnel; bands narrow from top to bottom.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, 1.30,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.32)
        write_paragraph(tb.text_frame, subtitle, size=typo.body_size,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.95
    body_bottom = layout.footer_top_in - 0.20
    body_h = body_bottom - body_top

    # Funnel occupies left half, descriptions on right
    funnel_left = layout.margin_left_in + 0.5
    funnel_right_max = layout.margin_left_in + width * 0.55
    funnel_w_top = funnel_right_max - funnel_left
    funnel_w_bot = funnel_w_top * 0.30
    funnel_cx = (funnel_left + funnel_right_max) / 2

    n = max(len(stages), 1)
    band_h = (body_h - 0.10 * (n - 1)) / n

    # Color gradient navy -> bright blue
    colors = [pal.deep_navy, pal.mid_blue, pal.bright_blue,
              pal.light_blue, pal.bright_blue, pal.light_blue]

    for i, st in enumerate(stages):
        # Width interpolation linear from full to bot
        t_top = i / n
        t_bot = (i + 1) / n
        w_top = funnel_w_top * (1 - t_top) + funnel_w_bot * t_top
        w_bot = funnel_w_top * (1 - t_bot) + funnel_w_bot * t_bot
        y_top = body_top + i * (band_h + 0.10)
        y_bot = y_top + band_h

        # Render as trapezoid via freeform on the slide
        # Approximation using a rectangle (simpler & still on-brand)
        rect_w = (w_top + w_bot) / 2
        rect_left = funnel_cx - rect_w / 2
        add_rect(slide, rect_left, y_top, rect_w, band_h,
                 fill=colors[i % len(colors)])

        # Stage name + value (white, centered)
        tb = add_textbox(slide, rect_left + 0.10, y_top, rect_w - 0.20,
                         band_h, anchor=MSO_ANCHOR.MIDDLE)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = st.get("name", f"Stage {i+1}")
        r1.font.size = Pt(typo.body_size + 1); r1.font.bold = True
        r1.font.color.rgb = pal.white; r1.font.name = typo.family
        if st.get("value"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = str(st["value"])
            r2.font.size = Pt(typo.body_size + 4); r2.font.bold = True
            r2.font.color.rgb = pal.white; r2.font.name = typo.family

        # Description on right
        right_left = funnel_right_max + 0.6
        right_w = layout.slide_width_in - layout.margin_right_in - right_left
        tb = add_textbox(slide, right_left, y_top, right_w, band_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, st.get("description", ""),
                        size=typo.body_size, color=pal.text_dark,
                        family=typo.family, first=True)
    return slide


# ---------- KPI dashboard ----------

def add_kpi_dashboard(prs, *,
                      title="[KPI dashboard / Insert action title]",
                      subtitle: Optional[str] = None,
                      kpis: Sequence[Dict],
                      columns: int = 4,
                      page_number=None,
                      section_marker="Section marker",
                      source="xx", footnote="1. xx",
                      theme: Theme = DEFAULT_THEME):
    """kpis: [{"label": "ARR", "value": "$1.2B", "delta": "+12% YoY",
                "delta_dir": "up"|"down"|"flat", "context": "..."}]
    Renders 4-8 tiles arranged in `columns` x ceil(n/columns) grid.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    if subtitle:
        tb = add_textbox(slide, layout.margin_left_in, 1.30,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.32)
        write_paragraph(tb.text_frame, subtitle, size=typo.body_size,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.95
    body_bottom = layout.footer_top_in - 0.20
    body_h = body_bottom - body_top

    n = max(len(kpis), 1)
    cols = max(1, min(columns, n))
    rows = (n + cols - 1) // cols
    gap = 0.20
    tile_w = (width - (cols - 1) * gap) / cols
    tile_h = (body_h - (rows - 1) * gap) / rows

    delta_color_map = {
        "up": pal.status_green,
        "down": pal.status_red,
        "flat": pal.footer_gray,
    }
    delta_glyph_map = {
        "up": "▲",
        "down": "▼",
        "flat": "▬",
    }

    for i, k in enumerate(kpis):
        ci = i % cols
        ri = i // cols
        x = layout.margin_left_in + ci * (tile_w + gap)
        y = body_top + ri * (tile_h + gap)
        # Tile background
        add_rect(slide, x, y, tile_w, tile_h, fill=pal.soft_gray)
        # Top accent bar
        add_rect(slide, x, y, tile_w, 0.10, fill=pal.bright_blue)

        # Label (top)
        tb = add_textbox(slide, x + 0.20, y + 0.20, tile_w - 0.40, 0.32)
        write_paragraph(tb.text_frame, k.get("label", "[KPI]"),
                        size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)

        # Big value (centered)
        val_top = y + 0.55
        val_h = tile_h - 1.20
        tb = add_textbox(slide, x + 0.20, val_top, tile_w - 0.40, val_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(k.get("value", "—")),
                        size=typo.title_size + 12, bold=True,
                        color=pal.deep_navy, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

        # Delta + glyph
        delta_dir = k.get("delta_dir", "flat")
        delta_text = k.get("delta")
        if delta_text:
            tb = add_textbox(slide, x + 0.20, y + tile_h - 0.55,
                             tile_w - 0.40, 0.30, anchor=MSO_ANCHOR.MIDDLE)
            tf = tb.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r1 = p.add_run()
            r1.text = delta_glyph_map.get(delta_dir, "") + " "
            r1.font.size = Pt(typo.body_size); r1.font.bold = True
            r1.font.color.rgb = delta_color_map.get(delta_dir, pal.footer_gray)
            r1.font.name = typo.family
            r2 = p.add_run()
            r2.text = str(delta_text)
            r2.font.size = Pt(typo.body_size); r2.font.bold = True
            r2.font.color.rgb = delta_color_map.get(delta_dir, pal.footer_gray)
            r2.font.name = typo.family

        # Context line (very bottom, small)
        if k.get("context"):
            tb = add_textbox(slide, x + 0.20, y + tile_h - 0.27,
                             tile_w - 0.40, 0.22)
            write_paragraph(tb.text_frame, k["context"],
                            size=typo.chart_label_size, color=pal.footer_gray,
                            family=typo.family, align=PP_ALIGN.CENTER,
                            first=True)
    return slide
