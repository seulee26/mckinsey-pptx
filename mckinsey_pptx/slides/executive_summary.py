"""Executive summary slides — paragraph variant + key-takeaway variant."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN

from ..base import (
    blank_slide, add_chrome, add_textbox, write_paragraph, add_subtitle_placeholder,
)
from ..theme import Theme, DEFAULT_THEME


def add_paragraph_summary(prs, *, title="Executive summary",
                           paragraphs: Sequence[str],
                           subtitle: Optional[str] = None,
                           page_number: Optional[int] = None,
                           section_marker: str = "Section marker",
                           source: Optional[str] = None,
                           footnote: Optional[str] = None,
                           theme: Theme = DEFAULT_THEME):
    """Variant A: paragraph-style summary (matches '1.png')."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)

    pal, typo, layout = theme.palette, theme.typography, theme.layout
    body_left = layout.margin_left_in
    body_top = layout.body_top_in
    body_w = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=body_top)
        body_top += 0.55

    tb = add_textbox(slide, body_left, body_top, body_w,
                     layout.footer_top_in - body_top - 0.1)
    first = True
    for para in paragraphs:
        write_paragraph(tb.text_frame, para, size=typo.body_size,
                        color=pal.text_dark, family=typo.family,
                        space_after=12, first=first)
        first = False
    return slide


def add_keytakeaway_summary(prs, *, title="Executive summary",
                            sections: Sequence[dict],
                            final_conclusion: Optional[str] = None,
                            page_number: Optional[int] = None,
                            section_marker: Optional[str] = None,
                            source: Optional[str] = None,
                            footnote: Optional[str] = None,
                            theme: Theme = DEFAULT_THEME):
    """Variant B: key-takeaway/bullets structure (matches '2.webp').

    sections: list of {"takeaway": str, "bullets": [str, ...]}
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)

    pal, typo, layout = theme.palette, theme.typography, theme.layout
    body_left = layout.margin_left_in
    body_top = layout.body_top_in
    body_w = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_h = layout.footer_top_in - body_top - 0.1

    tb = add_textbox(slide, body_left, body_top, body_w, body_h)
    first = True
    for i, sec in enumerate(sections):
        write_paragraph(tb.text_frame, sec["takeaway"], size=typo.body_size + 1,
                        bold=True, color=pal.text_dark, family=typo.family,
                        space_before=0 if first else 10, space_after=4,
                        first=first)
        first = False
        for b in sec.get("bullets", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, level=1)
    if final_conclusion:
        write_paragraph(tb.text_frame, final_conclusion,
                        size=typo.body_size + 1, bold=True,
                        color=pal.text_dark, family=typo.family,
                        space_before=14)
    return slide
