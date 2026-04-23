"""Trends / key areas slides — 4 variants."""
from __future__ import annotations
from typing import Sequence, Optional, Dict

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph, add_subtitle_placeholder,
)
from ..theme import Theme, DEFAULT_THEME


def _add_arrow_right(slide, theme, left, top, w=0.45, h=0.35):
    pal = theme.palette
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(left), Inches(top),
                               Inches(w), Inches(h))
    s.shadow.inherit = False
    s.fill.solid(); s.fill.fore_color.rgb = pal.bright_blue
    s.line.fill.background()
    return s


def _add_circle_number(slide, theme, left, top, d, n, fill=None, color=None):
    pal, typo = theme.palette, theme.typography
    fill = fill or pal.dark_navy
    color = color or pal.white
    add_oval(slide, left, top, d, d, fill=fill)
    tb = add_textbox(slide, left, top, d, d, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, str(n), size=typo.body_size, bold=True,
                    color=color, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)


# ---------- Three trends with icons ----------

def add_three_trends_icons(prs, *,
                           title="[Three key trends or areas slide / Insert action title]",
                           subtitle="[Insert subtitle]",
                           trends: Sequence[Dict],
                           page_number=None, section_marker="Section marker",
                           source="xx", footnote=None,
                           theme: Theme = DEFAULT_THEME):
    """trends: [{"label", "bullets": [str,...], "icon": optional}]"""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    body_top = 1.95
    block_h = (layout.footer_top_in - 0.2 - body_top) / max(len(trends), 1)
    icon_d = 0.85
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    for i, tr in enumerate(trends):
        y = body_top + i * block_h
        cy = y + block_h / 2
        # Circle icon
        add_oval(slide, layout.margin_left_in, cy - icon_d / 2,
                 icon_d, icon_d, fill=pal.deep_navy)
        # icon symbol (single Unicode char)
        ic = tr.get("icon", "★")
        tb = add_textbox(slide, layout.margin_left_in, cy - icon_d / 2,
                         icon_d, icon_d, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, ic, size=18, color=pal.white,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
        # Label
        text_left = layout.margin_left_in + icon_d + 0.25
        tb = add_textbox(slide, text_left, y + 0.05, width - icon_d - 0.25, 0.35)
        write_paragraph(tb.text_frame, f"[{tr.get('label','')}]",
                        size=typo.section_title_size, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
        # Bullets
        tb = add_textbox(slide, text_left, y + 0.40, width - icon_d - 0.25,
                         block_h - 0.40)
        first = True
        for b in tr.get("bullets", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
    return slide


# ---------- Three trends table (Trend / Description / Examples) ----------

def add_three_trends_table(prs, *,
                           title="[Three key trends slide / Insert action title]",
                           trends: Sequence[Dict],
                           page_number=None, section_marker="Section marker",
                           source="xx", footnote="1. xx",
                           theme: Theme = DEFAULT_THEME):
    """trends: [{"name", "description": [bullets], "examples": [bullets]}]"""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    body_top = 1.5
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    name_w = 1.7
    desc_w = 6.5
    exam_w = width - name_w - desc_w

    # Header row
    headers = [("Trend", layout.margin_left_in, name_w),
               ("Description", layout.margin_left_in + name_w + 0.2, desc_w),
               ("Examples", layout.margin_left_in + name_w + 0.2 + desc_w + 0.2,
                exam_w - 0.4)]
    for txt, x, w in headers:
        tb = add_textbox(slide, x, body_top, w, 0.32)
        write_paragraph(tb.text_frame, txt, size=typo.section_title_size,
                        bold=True, color=pal.text_dark, family=typo.family,
                        first=True)
    add_line(slide, layout.margin_left_in, body_top + 0.32,
             layout.margin_left_in + width, body_top + 0.32,
             color=pal.rule_gray, width_pt=0.5)

    rows_top = body_top + 0.45
    row_h = (layout.footer_top_in - 0.2 - rows_top) / max(len(trends), 1)

    for i, tr in enumerate(trends):
        ry = rows_top + i * row_h
        # Name pill
        add_rect(slide, layout.margin_left_in, ry + 0.05, name_w, row_h - 0.30,
                 fill=pal.dark_navy)
        tb = add_textbox(slide, layout.margin_left_in, ry + 0.05,
                         name_w, row_h - 0.30, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"[{tr.get('name','')}]",
                        size=typo.section_title_size, bold=True,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Description bullets
        tb = add_textbox(slide, layout.margin_left_in + name_w + 0.2, ry,
                         desc_w, row_h - 0.10)
        first = True
        for b in tr.get("description", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
        # Examples
        tb = add_textbox(slide, layout.margin_left_in + name_w + 0.2 + desc_w + 0.2,
                         ry, exam_w - 0.4, row_h - 0.10)
        first = True
        for b in tr.get("examples", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
        # Dashed bottom rule
        if i < len(trends) - 1:
            add_line(slide, layout.margin_left_in, ry + row_h - 0.05,
                     layout.margin_left_in + width, ry + row_h - 0.05,
                     color=pal.placeholder_gray, width_pt=0.5)
    return slide


# ---------- Three trends numbered ----------

def add_three_trends_numbered(prs, *,
                              title="[Three key trends slide / Insert action title]",
                              subtitle="[Insert subtitle]",
                              trends: Sequence[Dict],
                              page_number=None, section_marker="Section marker",
                              source="xx", footnote=None,
                              theme: Theme = DEFAULT_THEME):
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    body_top = 1.95
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    n = len(trends)
    block_h = (layout.footer_top_in - 0.2 - body_top) / max(n, 1)
    num_d = 0.55
    pill_w = 1.7
    text_left = layout.margin_left_in + num_d + 0.15 + pill_w + 0.3

    for i, tr in enumerate(trends):
        y = body_top + i * block_h
        cy = y + block_h / 2
        _add_circle_number(slide, theme, layout.margin_left_in, cy - num_d / 2,
                            num_d, i + 1)
        # Light blue pill with label
        add_rect(slide, layout.margin_left_in + num_d + 0.15,
                 y + 0.10, pill_w, block_h - 0.40, fill=pal.bright_blue)
        tb = add_textbox(slide, layout.margin_left_in + num_d + 0.15,
                         y + 0.10, pill_w, block_h - 0.40,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"[{tr.get('label','')}]",
                        size=typo.section_title_size, bold=True,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Bullets
        tb = add_textbox(slide, text_left, y + 0.10,
                         width - (text_left - layout.margin_left_in),
                         block_h - 0.20, anchor=MSO_ANCHOR.MIDDLE)
        first = True
        for b in tr.get("bullets", []):
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
    return slide


# ---------- Five key areas (numbered, with arrow + description) ----------

def add_five_key_areas(prs, *,
                       title="[Five key areas slide / Insert action title]",
                       subtitle="[Insert subtitle]",
                       areas: Sequence[Dict],
                       page_number=None, section_marker="Section marker",
                       source="xx", footnote=None,
                       theme: Theme = DEFAULT_THEME):
    """areas: [{"name", "description"}]"""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout
    if subtitle:
        add_subtitle_placeholder(slide, subtitle, theme, top_in=1.30)

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = 1.85
    # Header row
    name_col_x = layout.margin_left_in + 0.7
    arrow_x = layout.margin_left_in + 3.2
    desc_col_x = layout.margin_left_in + 4.4
    tb = add_textbox(slide, name_col_x, body_top, 3.0, 0.32)
    write_paragraph(tb.text_frame, "Area", size=typo.section_title_size,
                    bold=True, color=pal.text_dark, family=typo.family,
                    first=True)
    tb = add_textbox(slide, desc_col_x, body_top, width - 4.4, 0.32)
    write_paragraph(tb.text_frame, "Description", size=typo.section_title_size,
                    bold=True, color=pal.text_dark, family=typo.family,
                    first=True)
    add_line(slide, name_col_x, body_top + 0.32,
             layout.margin_left_in + width, body_top + 0.32,
             color=pal.rule_gray, width_pt=0.5)

    rows_top = body_top + 0.40
    n = len(areas)
    row_h = (layout.footer_top_in - 0.2 - rows_top) / max(n, 1)
    num_d = 0.45

    for i, area in enumerate(areas):
        ry = rows_top + i * row_h
        cy = ry + row_h / 2
        # Alternating row background
        if i % 2 == 0:
            add_rect(slide, layout.margin_left_in, ry, width, row_h,
                     fill=pal.soft_gray)
        _add_circle_number(slide, theme, layout.margin_left_in + 0.1,
                            cy - num_d / 2, num_d, i + 1)
        tb = add_textbox(slide, name_col_x, cy - 0.15, 3.0, 0.32,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"[{area.get('name','')}]",
                        size=typo.body_size, bold=True, color=pal.text_dark,
                        family=typo.family, first=True)
        # Arrow between cols (positioned safely between name and description)
        _add_arrow_right(slide, theme, arrow_x, cy - 0.14, 0.45, 0.28)
        # Description
        tb = add_textbox(slide, desc_col_x, cy - 0.18,
                         width - (desc_col_x - layout.margin_left_in), 0.4,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, area.get("description", ""),
                        size=typo.body_size, color=pal.text_dark,
                        family=typo.family, first=True)
    return slide
