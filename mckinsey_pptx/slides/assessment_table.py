"""Assessment / status overview table with traffic-light dots."""
from __future__ import annotations
from typing import Sequence, Optional, Literal

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_textbox,
    write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME

StatusColor = Literal["green", "amber", "red"]


def _status_rgb(theme: Theme, status: StatusColor):
    return {
        "green": theme.palette.status_green,
        "amber": theme.palette.status_amber,
        "red": theme.palette.status_red,
    }[status]


def add_assessment_table(prs, *, title="[Assessment or status overview / Insert action title]",
                         categories: Sequence[dict],
                         columns: Sequence[str] = ("Key Performance Indicators",
                                                    "Target", "Actual", "Status"),
                         page_number: Optional[int] = None,
                         section_marker: str = "Section marker",
                         source: Optional[str] = "xx",
                         footnote: Optional[str] = "1. xx",
                         theme: Theme = DEFAULT_THEME):
    """categories: [{"name": str, "rows": [{"kpi": str, "target": str,
                     "actual": str, "status_label": str, "status": "green"|"amber"|"red"}]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)

    pal, typo, layout = theme.palette, theme.typography, theme.layout
    left = layout.margin_left_in
    top = layout.body_top_in
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    # Column widths: Category | KPI (wider) | Target | Actual | Status
    cat_w = 1.6
    kpi_w = 4.4
    tgt_w = 1.3
    act_w = 1.3
    sts_w = width - cat_w - kpi_w - tgt_w - act_w
    col_xs = [left, left + cat_w, left + cat_w + kpi_w,
              left + cat_w + kpi_w + tgt_w,
              left + cat_w + kpi_w + tgt_w + act_w]
    col_ws = [cat_w, kpi_w, tgt_w, act_w, sts_w]

    # Header row
    header_h = 0.4
    header_titles = ["Category"] + list(columns)
    for x, w, txt in zip(col_xs, col_ws, header_titles):
        tb = add_textbox(slide, x + 0.1, top, w - 0.2, header_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, txt, size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER if txt != "Category" else PP_ALIGN.CENTER,
                        first=True)
    y = top + header_h

    # Body rows
    row_h = 0.4
    alt = True
    for cat in categories:
        rows = cat.get("rows", [])
        cat_h = row_h * len(rows) if rows else row_h
        # Category background (gray)
        add_rect(slide, col_xs[0], y, col_ws[0], cat_h, fill=pal.light_gray)
        cat_tb = add_textbox(slide, col_xs[0] + 0.1, y, col_ws[0] - 0.2, cat_h,
                             anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(cat_tb.text_frame, cat.get("name", ""),
                        size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

        for i, row in enumerate(rows):
            ry = y + i * row_h
            # alternating row backgrounds for the data area
            row_fill = pal.soft_gray if (i + (1 if alt else 0)) % 2 == 0 else None
            for cx, cw in zip(col_xs[1:], col_ws[1:]):
                add_rect(slide, cx, ry, cw, row_h,
                         fill=row_fill if row_fill else pal.white)

            # KPI text
            tb = add_textbox(slide, col_xs[1] + 0.1, ry, col_ws[1] - 0.2, row_h,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, row.get("kpi", ""),
                            size=typo.body_size, color=pal.text_dark,
                            family=typo.family, first=True)
            # Target
            tb = add_textbox(slide, col_xs[2], ry, col_ws[2], row_h,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, str(row.get("target", "")),
                            size=typo.body_size, color=pal.text_dark,
                            family=typo.family, align=PP_ALIGN.CENTER,
                            first=True)
            # Actual
            tb = add_textbox(slide, col_xs[3], ry, col_ws[3], row_h,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, str(row.get("actual", "")),
                            size=typo.body_size, color=pal.text_dark,
                            family=typo.family, align=PP_ALIGN.CENTER,
                            first=True)
            # Status label + dot
            tb = add_textbox(slide, col_xs[4] + 0.15, ry,
                             col_ws[4] - 0.6, row_h,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, str(row.get("status_label", "")),
                            size=typo.body_size, color=pal.text_dark,
                            family=typo.family, first=True)
            dot_d = 0.22
            dot_x = col_xs[4] + col_ws[4] - dot_d - 0.2
            dot_y = ry + (row_h - dot_d) / 2
            add_oval(slide, dot_x, dot_y, dot_d, dot_d,
                     fill=_status_rgb(theme, row.get("status", "green")))

        y += cat_h
        alt = not alt
    return slide
