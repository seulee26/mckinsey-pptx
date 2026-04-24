"""Structural slides every consulting deck needs:
 - cover_slide: title page with client / date / accent stripe
 - section_divider: chapter divider with big number + chapter name
 - agenda: numbered chapter list (with optional active highlight)
 - stat_hero: one big number + context label + supporting text
 - quote_slide: large pull-quote with attribution
"""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph, enable_text_shrink,
)
from ..theme import Theme, DEFAULT_THEME


# ---------- Cover slide ----------

def add_cover_slide(prs, *,
                    title: str,
                    subtitle: Optional[str] = None,
                    client: Optional[str] = None,
                    date: Optional[str] = None,
                    confidentiality: Optional[str] = "CONFIDENTIAL",
                    page_number=None,
                    section_marker=None,
                    source=None, footnote=None,
                    theme: Theme = DEFAULT_THEME):
    """First-page cover: large title, subtitle, client + date metadata.
    A 4-inch deep navy stripe runs down the right edge for visual weight.
    """
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Right-side accent stripe
    stripe_w = 0.40
    add_rect(slide, layout.slide_width_in - stripe_w, 0,
             stripe_w, layout.slide_height_in, fill=pal.deep_navy)

    # Top-right confidentiality tag (above stripe area)
    if confidentiality:
        tb = add_textbox(slide, layout.slide_width_in - 3.2,
                         0.30, 2.6, 0.25)
        write_paragraph(tb.text_frame, confidentiality,
                        size=typo.small_size, color=pal.placeholder_gray,
                        family=typo.family, align=PP_ALIGN.RIGHT, first=True)

    # Big title (centered vertically in upper half)
    title_left = layout.margin_left_in + 0.15
    title_top = 2.4
    title_w = layout.slide_width_in - layout.margin_left_in - 1.0
    tb = add_textbox(slide, title_left, title_top, title_w, 1.6)
    write_paragraph(tb.text_frame, title,
                    size=typo.title_size + 16, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    enable_text_shrink(tb.text_frame)

    # Subtitle
    if subtitle:
        tb = add_textbox(slide, title_left, title_top + 1.55, title_w, 0.80)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.title_size - 4,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)

    # Bottom line + client + date metadata
    add_line(slide, title_left, layout.slide_height_in - 1.2,
             title_left + 6.5, layout.slide_height_in - 1.2,
             color=pal.deep_navy, width_pt=1.0)
    meta_y = layout.slide_height_in - 1.05
    if client:
        tb = add_textbox(slide, title_left, meta_y, 6.0, 0.30)
        write_paragraph(tb.text_frame, client,
                        size=typo.body_size + 1, bold=True,
                        color=pal.text_dark, family=typo.family, first=True)
    if date:
        tb = add_textbox(slide, title_left, meta_y + 0.32, 6.0, 0.30)
        write_paragraph(tb.text_frame, date,
                        size=typo.body_size, color=pal.footer_gray,
                        family=typo.family, first=True)
    return slide


# ---------- Section divider ----------

def add_section_divider(prs, *,
                        section_number: str,
                        section_title: str,
                        subtitle: Optional[str] = None,
                        page_number=None,
                        section_marker=None,
                        source=None, footnote=None,
                        theme: Theme = DEFAULT_THEME):
    """Full-bleed section divider. Left third: deep-navy panel with the
    section number; right two-thirds: title + subtitle.
    """
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Left panel
    panel_w = 4.5
    add_rect(slide, 0, 0, panel_w, layout.slide_height_in, fill=pal.deep_navy)

    # Section number (very large, white, in left panel)
    tb = add_textbox(slide, 0.5, layout.slide_height_in / 2 - 1.5,
                     panel_w - 1.0, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, str(section_number),
                    size=typo.title_size + 56, bold=True,
                    color=pal.bright_blue, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    # Right side title
    right_left = panel_w + 0.6
    right_w = layout.slide_width_in - right_left - layout.margin_right_in
    tb = add_textbox(slide, right_left,
                     layout.slide_height_in / 2 - 1.0,
                     right_w, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, section_title,
                    size=typo.title_size + 8, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)

    # Underline accent
    add_line(slide, right_left,
             layout.slide_height_in / 2 + 0.40,
             right_left + 1.6,
             layout.slide_height_in / 2 + 0.40,
             color=pal.bright_blue, width_pt=2.5)

    if subtitle:
        tb = add_textbox(slide, right_left,
                         layout.slide_height_in / 2 + 0.55,
                         right_w, 1.0)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size + 2, color=pal.footer_gray,
                        family=typo.family, first=True)
    return slide


# ---------- Agenda ----------

def add_agenda(prs, *,
               title: str = "Agenda",
               items: Sequence[str],
               active_index: Optional[int] = None,
               page_number=None,
               section_marker=None,
               source=None, footnote=None,
               theme: Theme = DEFAULT_THEME):
    """Numbered agenda. Optional `active_index` highlights the current section
    in bright blue."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    body_top = 2.0
    body_bottom = layout.footer_top_in - 0.30
    n = max(len(items), 1)
    row_h = min(0.85, (body_bottom - body_top) / n)

    num_d = 0.55
    text_left_offset = num_d + 0.40

    for i, item in enumerate(items):
        y = body_top + i * row_h
        cy = y + row_h / 2
        is_active = (active_index is not None and i == active_index)
        circle_fill = pal.bright_blue if is_active else pal.deep_navy

        # Number circle
        add_oval(slide, layout.margin_left_in + 0.5,
                 cy - num_d / 2, num_d, num_d, fill=circle_fill)
        tb = add_textbox(slide, layout.margin_left_in + 0.5,
                         cy - num_d / 2, num_d, num_d,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(i + 1).zfill(2),
                        size=typo.body_size, bold=True, color=pal.white,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

        # Item label
        tb = add_textbox(slide,
                         layout.margin_left_in + 0.5 + text_left_offset,
                         cy - 0.25, 9.5, 0.50,
                         anchor=MSO_ANCHOR.MIDDLE)
        color = pal.bright_blue if is_active else pal.text_dark
        write_paragraph(tb.text_frame, item,
                        size=typo.body_size + 4, bold=is_active,
                        color=color, family=typo.family, first=True)

        # Subtle row divider
        if i < n - 1:
            add_line(slide,
                     layout.margin_left_in + 0.5,
                     y + row_h,
                     layout.slide_width_in - layout.margin_right_in - 0.5,
                     y + row_h,
                     color=pal.grid_gray, width_pt=0.4)
    return slide


# ---------- Stat hero ----------

def add_stat_hero(prs, *,
                  title: Optional[str] = None,
                  stat: str,
                  stat_label: str,
                  context: Optional[str] = None,
                  source_text: Optional[str] = None,
                  page_number=None,
                  section_marker="Section marker",
                  source=None, footnote=None,
                  theme: Theme = DEFAULT_THEME):
    """One enormous number + a one-line label + optional context paragraph.
    Use for "by 2030, 50% of all ..." style impact moments.
    """
    slide = blank_slide(prs)
    if title:
        add_chrome(slide, title=title, theme=theme, page_number=page_number,
                   section_marker=section_marker, source=source,
                   footnote=footnote)
        body_top = 2.0
    else:
        body_top = 1.0
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Big number, centered horizontally, in deep navy
    num_h = 2.6
    tb = add_textbox(slide, layout.margin_left_in,
                     body_top + 0.4, layout.slide_width_in
                     - layout.margin_left_in - layout.margin_right_in,
                     num_h, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, stat,
                    size=typo.title_size + 80, bold=True,
                    color=pal.deep_navy, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)
    enable_text_shrink(tb.text_frame)

    # Stat label below number
    tb = add_textbox(slide, layout.margin_left_in,
                     body_top + 0.4 + num_h + 0.05,
                     layout.slide_width_in - layout.margin_left_in
                     - layout.margin_right_in, 0.50)
    write_paragraph(tb.text_frame, stat_label,
                    size=typo.title_size, bold=True,
                    color=pal.bright_blue, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    # Optional context paragraph
    if context:
        tb = add_textbox(slide, 2.0, body_top + 0.4 + num_h + 0.85,
                         layout.slide_width_in - 4.0, 1.2)
        write_paragraph(tb.text_frame, context, size=typo.body_size + 2,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

    if source_text and not title:
        # Source line at bottom (when no chrome)
        tb = add_textbox(slide, layout.margin_left_in,
                         layout.slide_height_in - 0.45,
                         layout.slide_width_in - layout.margin_left_in
                         - layout.margin_right_in, 0.25)
        write_paragraph(tb.text_frame, f"Source: {source_text}",
                        size=typo.footer_size, color=pal.footer_gray,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)
    return slide


# ---------- Quote slide ----------

def add_quote_slide(prs, *,
                    quote: str,
                    author: str,
                    author_title: Optional[str] = None,
                    title: Optional[str] = None,
                    page_number=None,
                    section_marker=None,
                    source=None, footnote=None,
                    theme: Theme = DEFAULT_THEME):
    """Pull-quote slide: oversized quotation marks, the quote in italic, and
    attribution lines below.
    """
    slide = blank_slide(prs)
    if title:
        add_chrome(slide, title=title, theme=theme, page_number=page_number,
                   section_marker=section_marker, source=source,
                   footnote=footnote)
        body_top = 1.95
    else:
        body_top = 1.5
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Big opening quotation mark
    tb = add_textbox(slide, 0.8, body_top, 1.5, 1.5)
    write_paragraph(tb.text_frame, "“",
                    size=typo.title_size + 60, bold=True,
                    color=pal.bright_blue, family=typo.family, first=True)

    # Quote body
    quote_left = 2.0
    quote_w = layout.slide_width_in - quote_left - layout.margin_right_in - 0.5
    tb = add_textbox(slide, quote_left, body_top + 0.4, quote_w, 3.5)
    write_paragraph(tb.text_frame, quote,
                    size=typo.title_size, bold=False, italic=False,
                    color=pal.text_dark, family=typo.family, first=True)

    # Attribution
    attr_y = body_top + 4.0
    add_line(slide, quote_left, attr_y, quote_left + 0.6, attr_y,
             color=pal.bright_blue, width_pt=2.0)

    tb = add_textbox(slide, quote_left, attr_y + 0.10, quote_w, 0.35)
    write_paragraph(tb.text_frame, author, size=typo.body_size + 2,
                    bold=True, color=pal.text_dark, family=typo.family,
                    first=True)
    if author_title:
        tb = add_textbox(slide, quote_left, attr_y + 0.45, quote_w, 0.35)
        write_paragraph(tb.text_frame, author_title,
                        size=typo.body_size, color=pal.footer_gray,
                        family=typo.family, first=True)
    return slide
