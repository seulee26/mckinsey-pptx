"""Bubble / scatter slides — 4 variants:
 - bubble chart (full width)
 - bubble chart with takeaways (chart left, bullets right)
 - growth-share matrix (BCG style 2x2)
 - prioritization / assessment matrix (3x3 with status colors)
"""
from __future__ import annotations
from typing import Sequence, Optional, Tuple, List, Dict, Literal
import math

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_oval, add_line, add_textbox,
    write_paragraph,
)
from ..theme import Theme, DEFAULT_THEME
from .column_chart import _draw_takeaway, _draw_description_header


def _round_up_nice(v):
    if v <= 0:
        return 1
    mag = 10 ** math.floor(math.log10(v))
    n = v / mag
    if n <= 1: pick = 1
    elif n <= 2: pick = 2
    elif n <= 2.5: pick = 2.5
    elif n <= 5: pick = 5
    else: pick = 10
    return pick * mag


def _draw_xy_axis(slide, theme, *, plot_box, x_max, y_max,
                  x_label, x_unit, y_label, y_unit, n_ticks_x=10, n_ticks_y=7):
    pal, typo = theme.palette, theme.typography
    pl, pt, pw, ph = plot_box
    pr = pl + pw
    pb = pt + ph

    x_ticks = [i * x_max / (n_ticks_x - 1) for i in range(n_ticks_x)]
    y_ticks = [i * y_max / (n_ticks_y - 1) for i in range(n_ticks_y)]

    # Grid
    for tv in y_ticks:
        ty = pb - (tv / y_max) * ph
        add_line(slide, pl, ty, pr, ty, color=pal.grid_gray, width_pt=0.5)
    for tv in x_ticks:
        tx = pl + (tv / x_max) * pw
        add_line(slide, tx, pt, tx, pb, color=pal.grid_gray, width_pt=0.5)

    # Axis labels (numbers)
    for tv in y_ticks:
        ty = pb - (tv / y_max) * ph
        tb = add_textbox(slide, pl - 0.55, ty - 0.12, 0.45, 0.24,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"{int(round(tv))}",
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.RIGHT, first=True)
    for tv in x_ticks:
        tx = pl + (tv / x_max) * pw
        tb = add_textbox(slide, tx - 0.30, pb + 0.05, 0.6, 0.22)
        write_paragraph(tb.text_frame, f"{int(round(tv))}",
                        size=typo.chart_axis_size, color=pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    # Axis titles
    tb = add_textbox(slide, pl - 0.6, pt - 0.45, 3.5, 0.3)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = f"[{y_label}], "
    r.font.size = Pt(typo.section_title_size); r.font.bold = True
    r.font.color.rgb = pal.text_dark; r.font.name = typo.family
    r2 = p.add_run(); r2.text = f"[{y_unit}]"
    r2.font.size = Pt(typo.section_title_size)
    r2.font.color.rgb = pal.placeholder_gray; r2.font.name = typo.family

    tb = add_textbox(slide, pr - 2.5, pb + 0.30, 2.5, 0.3)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"[{x_label}], "
    r.font.size = Pt(typo.section_title_size); r.font.bold = True
    r.font.color.rgb = pal.text_dark; r.font.name = typo.family
    r2 = p.add_run(); r2.text = f"[{x_unit}]"
    r2.font.size = Pt(typo.section_title_size)
    r2.font.color.rgb = pal.placeholder_gray; r2.font.name = typo.family


def _draw_bubble(slide, theme, *, plot_box, x_max, y_max, bubble,
                 size_min=0.25, size_max=0.95):
    pal = theme.palette
    pl, pt, pw, ph = plot_box
    pb = pt + ph

    sz_vals = [b.get("size", 1) for b in bubble]
    s_max = max(sz_vals) if sz_vals else 1
    s_min = min(sz_vals) if sz_vals else 1

    color_map = {
        "blue_light": pal.bright_blue,
        "blue_dark": pal.dark_navy,
        "blue_royal": pal.royal_blue,
        "navy": pal.dark_navy,
    }

    for b in bubble:
        x = pl + (b["x"] / x_max) * pw
        y = pb - (b["y"] / y_max) * ph
        s = b.get("size", 1)
        if s_max == s_min:
            d = (size_min + size_max) / 2
        else:
            d = size_min + (size_max - size_min) * (s - s_min) / (s_max - s_min)
        fill = color_map.get(b.get("group"), pal.bright_blue)
        if isinstance(b.get("color"), str):
            fill = color_map.get(b["color"], fill)
        add_oval(slide, x - d / 2, y - d / 2, d, d, fill=fill)
        # Label to the right of bubble
        if b.get("label"):
            tb = add_textbox(slide, x + d / 2 + 0.05, y - 0.12,
                              1.2, 0.24,
                              anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, b["label"],
                            size=theme.typography.chart_label_size,
                            color=pal.text_dark,
                            family=theme.typography.family, first=True)


def _draw_legend_groups(slide, theme, *, top_left, groups, with_size_swatch=True):
    """groups: [(color_name, label)]"""
    pal, typo = theme.palette, theme.typography
    color_map = {
        "blue_light": pal.bright_blue,
        "blue_dark": pal.dark_navy,
        "blue_royal": pal.royal_blue,
        "navy": pal.dark_navy,
    }
    x, y = top_left
    for cname, label in groups:
        rgb = color_map.get(cname, pal.bright_blue)
        d = 0.28
        add_oval(slide, x, y - 0.05, d, d, fill=rgb)
        tb = add_textbox(slide, x + d + 0.08, y - 0.06, 1.5, 0.30,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"[{label}]", size=typo.chart_label_size,
                        color=pal.text_dark, family=typo.family, first=True)
        x += d + 1.5
    if with_size_swatch:
        d = 0.28
        add_oval(slide, x, y - 0.05, d, d, fill=None,
                 line=pal.placeholder_gray, line_width=0.75)
        tb = add_textbox(slide, x + d + 0.08, y - 0.06, 2.4, 0.30,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, "Size = [Insert dimension]",
                        size=typo.chart_label_size, color=pal.text_dark,
                        family=typo.family, first=True)


# ---------- Public builders ----------

def add_bubble_chart(prs, *, title="[Bubble chart / Insert action title]",
                    bubbles: Sequence[Dict],
                    x_max=900, y_max=3000,
                    x_label="Dimension 2", x_unit="Unit",
                    y_label="Dimension 1", y_unit="Unit",
                    groups: Sequence[Tuple[str, str]] = (
                        ("blue_light", "Insert group 1"),
                        ("blue_dark", "Insert group 2"),
                        ("blue_royal", "Insert group 3"),
                    ),
                    diagonal=True,
                    state_top_left: Optional[str] = "State e.g. profitable",
                    state_bottom_right: Optional[str] = "State, e.g. unprofitable",
                    page_number=None, section_marker=None,
                    source="xx", footnote="1. xx",
                    theme: Theme = DEFAULT_THEME):
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Legend top
    _draw_legend_groups(slide, theme, top_left=(3.4, 1.55), groups=groups)

    plot_box = (1.05, 2.05, 11.7, 4.6)
    _draw_xy_axis(slide, theme, plot_box=plot_box, x_max=x_max, y_max=y_max,
                  x_label=x_label, x_unit=x_unit,
                  y_label=y_label, y_unit=y_unit)
    pl, pt, pw, ph = plot_box

    # Diagonal reference
    if diagonal:
        line = add_line(slide, pl, pt + ph, pl + pw, pt,
                         color=pal.text_dark, width_pt=0.75,
                         dash=MSO_LINE_DASH_STYLE.DASH)

    # State labels
    if state_top_left:
        tb = add_textbox(slide, pl + 0.3, pt + 0.15, 3.5, 0.3)
        write_paragraph(tb.text_frame, f"[{state_top_left}]",
                        size=typo.chart_label_size, italic=True,
                        color=pal.placeholder_gray, family=typo.family,
                        first=True)
    if state_bottom_right:
        tb = add_textbox(slide, pl + pw - 3.5, pt + ph - 0.45, 3.4, 0.3)
        write_paragraph(tb.text_frame, f"[{state_bottom_right}]",
                        size=typo.chart_label_size, italic=True,
                        color=pal.placeholder_gray, family=typo.family,
                        align=PP_ALIGN.RIGHT, first=True)

    _draw_bubble(slide, theme, plot_box=plot_box, x_max=x_max, y_max=y_max,
                 bubble=bubbles)
    return slide


def add_bubble_chart_with_takeaways(prs, *,
                                    title="[Bubble chart with takeaways / Insert action title]",
                                    bubbles, x_max=900, y_max=3000,
                                    x_label="Dimension 2", x_unit="Unit",
                                    y_label="Dimension 1", y_unit="Unit",
                                    groups=(
                                        ("blue_dark", "Insert group 1"),
                                        ("blue_light", "Insert group 2"),
                                        ("mid_blue", "Insert group 3"),
                                    ),
                                    takeaways=(),
                                    page_number=None, section_marker=None,
                                    source="xx", footnote=None,
                                    theme: Theme = DEFAULT_THEME):
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    # Description header on left side
    _draw_description_header(slide, theme, left=0.45, top=1.45, width=8.5)

    plot_box = (1.05, 2.30, 7.9, 3.55)
    _draw_xy_axis(slide, theme, plot_box=plot_box, x_max=x_max, y_max=y_max,
                  x_label=x_label, x_unit=x_unit,
                  y_label=y_label, y_unit=y_unit)
    _draw_bubble(slide, theme, plot_box=plot_box, x_max=x_max, y_max=y_max,
                 bubble=bubbles)
    # Legend below chart (kept clear of axis title)
    _draw_legend_groups(slide, theme, top_left=(1.05, 6.55), groups=groups)
    # Takeaway right side
    _draw_takeaway(slide, theme, takeaways=takeaways,
                   box=(9.45, 1.45, 3.45, 5.0))
    return slide


def add_growth_share_matrix(prs, *,
                            title="[Growth-share matrix / Insert action title]",
                            bus: Sequence[Dict],
                            x_max=100, y_max=50,
                            page_number=None, section_marker=None,
                            source="xx", footnote="1. xx",
                            theme: Theme = DEFAULT_THEME):
    """bus: list of {"name", "x" (market share), "y" (growth), "size", "quadrant"?}"""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    plot_box = (1.05, 1.65, 11.7, 5.0)
    pl, pt, pw, ph = plot_box

    # Quadrant fills (split at 50% / 25)
    midx = pl + pw / 2
    midy = pt + ph / 2
    add_rect(slide, pl, pt, pw / 2, ph / 2, fill=pal.mid_blue)             # top-left
    add_rect(slide, midx, pt, pw / 2, ph / 2, fill=pal.bright_blue)         # top-right (Star)
    add_rect(slide, pl, midy, pw / 2, ph / 2, fill=pal.soft_gray)           # bottom-left (Dog)
    add_rect(slide, midx, midy, pw / 2, ph / 2, fill=pal.mid_blue)          # bottom-right

    # Quadrant labels
    label_pads = [
        ("Question mark", pl + 0.2, pt + 0.15, pal.white),
        ("Star", midx + 0.2, pt + 0.15, pal.white),
        ("Dog", pl + 0.2, midy + 0.15, pal.text_dark),
        ("Cash cow", midx + 0.2, midy + 0.15, pal.white),
    ]
    for label, lx, ly, color in label_pads:
        tb = add_textbox(slide, lx, ly, 2.5, 0.3)
        write_paragraph(tb.text_frame, label, size=typo.section_title_size,
                        color=color, family=typo.family, first=True)

    # Y-axis label
    tb = add_textbox(slide, pl - 0.55, pt - 0.45, 4, 0.3)
    write_paragraph(tb.text_frame, "Growth rate 20xx-20xx (%)",
                    size=typo.section_title_size, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    # X-axis label
    tb = add_textbox(slide, pl, pt + ph + 0.30, pw, 0.3)
    write_paragraph(tb.text_frame, "Market share (%)",
                    size=typo.section_title_size, bold=True,
                    color=pal.text_dark, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    # Y ticks
    for v in range(0, int(y_max) + 1, 5):
        ty = pt + ph - (v / y_max) * ph
        tb = add_textbox(slide, pl - 0.55, ty - 0.10, 0.45, 0.22,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(v), size=typo.chart_axis_size,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.RIGHT, first=True)
    for v in range(0, int(x_max) + 1, 10):
        tx = pl + (v / x_max) * pw
        tb = add_textbox(slide, tx - 0.30, pt + ph + 0.05, 0.6, 0.22)
        write_paragraph(tb.text_frame, str(v), size=typo.chart_axis_size,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

    # Size legend
    add_oval(slide, pl + pw - 3.0, pt - 0.45, 0.28, 0.28, fill=pal.dark_navy)
    tb = add_textbox(slide, pl + pw - 2.65, pt - 0.45, 2.6, 0.3,
                     anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, "Size = [insert description]",
                    size=typo.chart_label_size, color=pal.text_dark,
                    family=typo.family, first=True)

    # Bubbles: dark navy with thin white outline
    sz_vals = [b.get("size", 1) for b in bus]
    s_max = max(sz_vals) if sz_vals else 1
    s_min = min(sz_vals) if sz_vals else 1
    for b in bus:
        x = pl + (b["x"] / x_max) * pw
        y = pt + ph - (b["y"] / y_max) * ph
        s = b.get("size", 1)
        if s_max == s_min:
            d = 0.6
        else:
            d = 0.4 + 0.6 * (s - s_min) / (s_max - s_min)
        add_oval(slide, x - d / 2, y - d / 2, d, d, fill=pal.deep_navy,
                 line=pal.white, line_width=1.0)
        if b.get("name"):
            tb = add_textbox(slide, x + d / 2 + 0.06, y - 0.12, 1.0, 0.24,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, b["name"],
                            size=typo.chart_label_size, color=pal.white,
                            family=typo.family, first=True)
    return slide


def add_prioritization_matrix(prs, *,
                              title="[Prioritization or assessment matrix / Insert action title]",
                              items: Sequence[Dict],
                              page_number=None, section_marker=None,
                              source="xx", footnote="1. xx",
                              theme: Theme = DEFAULT_THEME):
    """items: [{"name", "x_band": 0|1|2 (Low/Med/High), "y_band": 0|1|2 (Long/Med/Short), "status": green|amber|red}]"""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number,
               section_marker=section_marker, source=source, footnote=footnote)
    pal, typo = theme.palette, theme.typography

    # Description left
    tb = add_textbox(slide, 0.45, 1.5, 4, 0.3)
    write_paragraph(tb.text_frame, "[Description]",
                    size=typo.section_title_size, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)

    # Legend top
    leg_items = [
        ("green", "Insert status/group"),
        ("amber", "Insert status/group"),
        ("red", "Insert status/group"),
    ]
    leg_x = 5.4
    leg_y = 1.5
    color_map = {
        "green": pal.status_green,
        "amber": pal.status_amber,
        "red": pal.status_red,
    }
    for k, label in leg_items:
        d = 0.28
        add_oval(slide, leg_x, leg_y, d, d, fill=color_map[k])
        tb = add_textbox(slide, leg_x + d + 0.08, leg_y - 0.02, 2.5, 0.32,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, f"= [{label}]",
                        size=typo.chart_label_size, color=pal.text_dark,
                        family=typo.family, first=True)
        leg_x += 2.5

    # Plot grid box (kept short of footer to leave room for axis title)
    plot_box = (1.55, 2.05, 11.30, 4.40)
    pl, pt, pw, ph = plot_box

    # Highlight top-right cell (Short × High)
    cell_w, cell_h = pw / 3, ph / 3
    add_rect(slide, pl + 2 * cell_w, pt, cell_w, cell_h, fill=pal.bright_blue)
    # Middle row gray strip
    add_rect(slide, pl + cell_w, pt + cell_h, cell_w * 2, cell_h,
             fill=pal.soft_gray)

    # Dashed grid lines
    for i in range(1, 3):
        gx = pl + i * cell_w
        add_line(slide, gx, pt, gx, pt + ph, color=pal.placeholder_gray,
                 width_pt=0.5, dash=MSO_LINE_DASH_STYLE.DASH)
        gy = pt + i * cell_h
        add_line(slide, pl, gy, pl + pw, gy, color=pal.placeholder_gray,
                 width_pt=0.5, dash=MSO_LINE_DASH_STYLE.DASH)

    # Y axis labels (Short / Medium / Long)
    for i, label in enumerate(["Short", "Medium", "Long"]):
        ly = pt + i * cell_h + 0.05
        tb = add_textbox(slide, pl - 0.95, ly, 0.85, 0.28)
        write_paragraph(tb.text_frame, label, size=typo.chart_label_size,
                        color=pal.placeholder_gray, family=typo.family,
                        align=PP_ALIGN.RIGHT, first=True)
    # X axis labels (Low / Medium / High)
    for i, label in enumerate(["Low", "Medium", "High"]):
        lx = pl + i * cell_w + cell_w / 2 - 0.5
        tb = add_textbox(slide, lx, pt + ph + 0.06, 1.0, 0.28)
        write_paragraph(tb.text_frame, label, size=typo.chart_label_size,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

    # Axis titles
    tb = add_textbox(slide, pl - 0.95, pt + ph / 2 - 0.5, 0.85, 1.0,
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "TIME TO IMPACT"
    r.font.size = Pt(typo.section_title_size); r.font.bold = True
    r.font.color.rgb = pal.text_dark; r.font.name = typo.family

    tb = add_textbox(slide, pl + pw - 2.5, pt + ph + 0.30, 2.5, 0.3)
    write_paragraph(tb.text_frame, "LEVEL OF IMPACT",
                    size=typo.section_title_size, bold=True,
                    color=pal.text_dark, family=typo.family,
                    align=PP_ALIGN.RIGHT, first=True)

    # Items
    for it in items:
        bx = it.get("x_band", 1)
        by = it.get("y_band", 1)  # 0=Short(top), 1=Medium, 2=Long(bottom)
        # offsets within cell
        ox = it.get("ox", 0.5)
        oy = it.get("oy", 0.5)
        x = pl + (bx + ox) * cell_w
        y = pt + (by + oy) * cell_h
        d = it.get("d", 0.85)
        add_oval(slide, x - d / 2, y - d / 2, d, d,
                 fill=color_map.get(it.get("status", "green"), pal.status_green))
        # Centered label
        tb = add_textbox(slide, x - d / 2, y - d / 2, d, d,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, it.get("name", ""),
                        size=typo.chart_label_size, color=pal.white
                        if it.get("status") != "amber" else pal.text_dark,
                        family=typo.family, align=PP_ALIGN.CENTER,
                        bold=False, first=True)
    return slide
